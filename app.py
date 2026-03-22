"""
AI Hub Web App (Cloud Deployment Version)
==========================================
Includes login authentication for security.

Environment Variables (set in Render Dashboard):
  OPENAI_API_KEY       - ChatGPT API key
  GEMINI_API_KEY       - Gemini API key
  AZURE_OPENAI_API_KEY - Azure OpenAI API key
  AZURE_OPENAI_ENDPOINT - Azure endpoint URL
  CLAUDE_API_KEY       - Claude (Anthropic) API key
  GROK_API_KEY         - Grok (xAI) API key
  APP_USERNAME         - Login username (default: admin)
  APP_PASSWORD         - Login password (default: aihub2026)
  SECRET_KEY           - Flask session secret (auto-generated if not set)
  SUPABASE_URL         - Supabase project URL
  SUPABASE_KEY         - Supabase anon public key
"""

import os
import tempfile
import secrets
import hashlib
from functools import wraps
from datetime import datetime, timedelta
from collections import defaultdict
import time

from flask import Flask, request, jsonify, session, redirect, url_for, Response, stream_with_context, render_template, send_from_directory
import json

from ai_hub import AIHub

app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY", secrets.token_hex(32))
app.permanent_session_lifetime = timedelta(hours=int(os.getenv("SESSION_TIMEOUT_HOURS", "2")))
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB upload limit

hub = AIHub()

APP_USERNAME = os.getenv("APP_USERNAME", "admin")
_raw_password = os.getenv("APP_PASSWORD", "aihub2026")

# Password hashing: SHA-256 with salt
def _hash_password(pw: str) -> str:
    salt = os.getenv("PASSWORD_SALT", "aihub_salt_2026")
    return hashlib.sha256(f"{salt}:{pw}".encode()).hexdigest()

APP_PASSWORD_HASH = _hash_password(_raw_password)


# ──────────────────────────── Rate Limiting (Tiered) ────────────────────────

class RateLimiter:
    """In-memory tiered rate limiter per IP."""
    TIERS = {
        "admin":   None,                               # unlimited
        "premium": {"requests": 60,  "window": 60},   # 60 req/min
        "free":    {"requests": 20,  "window": 60},   # 20 req/min
    }

    def __init__(self):
        self.requests = defaultdict(list)  # ip -> [timestamps]

    def is_allowed(self, ip: str, tier: str = "free") -> bool:
        cfg = self.TIERS.get(tier, self.TIERS["free"])
        if cfg is None:
            return True  # admin: unlimited
        now = time.time()
        self.requests[ip] = [t for t in self.requests[ip] if now - t < cfg["window"]]
        if len(self.requests[ip]) >= cfg["requests"]:
            return False
        self.requests[ip].append(now)
        return True

    def remaining(self, ip: str, tier: str = "free") -> int:
        cfg = self.TIERS.get(tier, self.TIERS["free"])
        if cfg is None:
            return -1  # unlimited
        now = time.time()
        self.requests[ip] = [t for t in self.requests[ip] if now - t < cfg["window"]]
        return max(0, cfg["requests"] - len(self.requests[ip]))

    def tier_info(self, tier: str = "free") -> dict:
        cfg = self.TIERS.get(tier, self.TIERS["free"])
        if cfg is None:
            return {"tier": tier, "max_requests": "unlimited", "window_seconds": 0}
        return {"tier": tier, "max_requests": cfg["requests"], "window_seconds": cfg["window"]}


api_limiter = RateLimiter()
login_limiter = RateLimiter()  # uses 'free' tier (20 attempts/min) for login

UPLOAD_DIR = tempfile.mkdtemp(prefix="aihub_uploads_")

# ──────────────────────────── Supabase ────────────────────────────

supabase_client = None
supabase_admin = None  # service_role client for admin ops (bypasses RLS)
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")

if SUPABASE_URL and SUPABASE_KEY:
    try:
        from supabase import create_client
        supabase_client = create_client(SUPABASE_URL, SUPABASE_KEY)
        print("  ✅ Supabase connected")
        if SUPABASE_SERVICE_KEY:
            supabase_admin = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
            print("  ✅ Supabase admin (service_role) connected")
        else:
            supabase_admin = supabase_client  # fallback
            print("  ⚠️ No SUPABASE_SERVICE_KEY, admin ops use anon key")
    except Exception as e:
        print(f"  ⚠️ Supabase init failed: {e}")
else:
    print("  ⚠️ Supabase not configured (no SUPABASE_URL/KEY)")


# ──────────────────────────── Authentication ────────────────────────────

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get("logged_in"):
            if request.path.startswith("/api/") or request.is_json:
                return jsonify({"success": False, "error": "Session expired. Please refresh and login again."}), 401
            return redirect(url_for("login_page"))
        # Session timeout check
        last_active = session.get("last_active")
        if last_active:
            try:
                elapsed = (datetime.utcnow() - datetime.fromisoformat(last_active)).total_seconds()
                if elapsed > app.permanent_session_lifetime.total_seconds():
                    session.clear()
                    if request.path.startswith("/api/"):
                        return jsonify({"success": False, "error": "Session timed out. Please login again."}), 401
                    return redirect(url_for("login_page"))
            except Exception:
                pass
        session["last_active"] = datetime.utcnow().isoformat()
        # Tiered rate limiting for API calls
        if request.path.startswith("/api/"):
            ip = request.remote_addr or "unknown"
            tier = session.get("user_tier", os.getenv("USER_TIER", "owner"))
            if not api_limiter.is_allowed(ip, tier):
                info = api_limiter.tier_info(tier)
                return jsonify({
                    "success": False,
                    "error": f"Rate limit exceeded ({info['max_requests']} req/{info['window_seconds']}s for {tier} tier). Try again shortly.",
                    "tier": tier,
                    "limit": info["max_requests"]
                }), 429
        return f(*args, **kwargs)
    return decorated


# PWA: Serve service worker from root scope
@app.route("/service-worker.js")
def service_worker():
    return send_from_directory(app.static_folder, "service-worker.js", mimetype="application/javascript")


@app.route("/login", methods=["GET", "POST"])
def login_page():
    if request.method == "POST":
        ip = request.remote_addr or "unknown"
        if not login_limiter.is_allowed(ip):
            return render_template("login.html", error_msg="Too many login attempts. Please wait and try again.")
        username = request.form.get("username", "")
        password = request.form.get("password", "")
        pw_hash = _hash_password(password)
        # Try Supabase users table first
        user_found = False
        db_client = supabase_admin or supabase_client
        if db_client:
            for _attempt in range(4):
                try:
                    res = db_client.table("users").select("*").eq("username", username).eq("is_active", True).execute()
                    if res.data and len(res.data) > 0:
                        db_user = res.data[0]
                        if db_user["password_hash"] == pw_hash:
                            session.permanent = True
                            session["logged_in"] = True
                            session["user_id"] = db_user["id"]
                            session["username"] = db_user["username"]
                            session["user_tier"] = db_user.get("tier", "free")
                            session["display_name"] = db_user.get("display_name", username)
                            session["last_active"] = datetime.utcnow().isoformat()
                            session["login_time"] = datetime.utcnow().isoformat()
                            session["must_change_password"] = bool(db_user.get("must_change_password"))
                            try:
                                db_client.table("users").update({"last_login": datetime.utcnow().isoformat()}).eq("id", db_user["id"]).execute()
                            except Exception:
                                pass
                            return redirect("/")
                        else:
                            user_found = True  # user exists but wrong password
                            print(f"  ⚠️ Login failed: user '{username}' found but password mismatch")
                    else:
                        print(f"  ⚠️ Login: user '{username}' not found in Supabase (is_active=True)")
                    break  # query succeeded, no need to retry
                except Exception as e:
                    err = str(e)
                    print(f"  ⚠️ Login Supabase error (attempt {_attempt+1}): {err}")
                    if ("PGRST002" in err or "schema cache" in err) and _attempt < 3:
                        time.sleep(1.0)
                    else:
                        break  # non-retryable error, fall through
        else:
            print("  ⚠️ Login: No Supabase client available, trying env fallback only")
        # Fallback: env var credentials
        if not user_found and username == APP_USERNAME and pw_hash == APP_PASSWORD_HASH:
            session.permanent = True
            session["logged_in"] = True
            session["user_id"] = "env_admin"
            session["username"] = username
            session["user_tier"] = "owner"
            session["display_name"] = "Owner"
            session["last_active"] = datetime.utcnow().isoformat()
            session["login_time"] = datetime.utcnow().isoformat()
            return redirect("/")
        return render_template("login.html", error_msg="Invalid credentials")
    return render_template("login.html")


def _seed_admin_user():
    """Auto-create admin user in Supabase if none exists."""
    if not supabase_client:
        return
    try:
        res = supabase_admin.table("users").select("id").eq("tier", "owner").limit(1).execute()
        if not res.data:
            supabase_admin.table("users").insert({
                "username": APP_USERNAME,
                "password_hash": APP_PASSWORD_HASH,
                "tier": "owner",
                "display_name": "System Owner",
                "is_active": True,
            }).execute()
            print("  ✅ Admin user seeded in Supabase")
            
        # Ensure shinwookyi is upgraded to owner
        try:
            supabase_admin.table("users").update({"tier": "owner"}).eq("username", "shinwookyi").execute()
            print("  ✅ Ensured shinwookyi is owner")
        except Exception as e:
            print(f"  ⚠️ Failed to force-upgrade shinwookyi: {e}")

        # Seed guest user (ensure password_hash matches current SALT)
        try:
            guest_hash = _hash_password("Guest123")
            existing = supabase_client.table("users").select("id").eq("username", "guest").execute()
            if existing.data:
                supabase_admin.table("users").update({"password_hash": guest_hash, "is_active": True}).eq("username", "guest").execute()
                print("  ✅ Guest user password synced")
            else:
                supabase_admin.table("users").insert({
                    "username": "guest", "password_hash": guest_hash,
                    "tier": "free", "display_name": "Guest", "is_active": True,
                }).execute()
                print("  ✅ Guest user created")
        except Exception as e:
            print(f"  ⚠️ Guest seed skipped: {e}")

    except Exception as e:
        print(f"  ⚠️ Admin seed skipped: {e}")


# Admin-only decorator
def admin_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if session.get("user_tier") != "owner":
            return jsonify({"success": False, "error": "Owner access required"}), 403
        return f(*args, **kwargs)
    return decorated


@app.route("/logout")
def logout():
    user_tier = session.get("user_tier", "free")
    username = session.get("username", "")
    
    # Guest/free tier: clear ALL data on logout
    if user_tier in ("free", "guest") or username == "guest":
        # Clear AI Hub chat history (server-side memory)
        hub.clear_history()
        
        # Delete conversations from Supabase
        try:
            if username:
                supabase_admin.table("conversations").delete().eq("username", username).execute()
        except Exception:
            pass
        
        # Delete persona memories from Supabase
        try:
            if username:
                supabase_admin.table("persona_memory").delete().eq("username", username).execute()
        except Exception:
            pass
    
    session.clear()
    # Redirect with cleanup flag for guest users
    if user_tier in ("free", "guest") or username == "guest":
        return redirect(url_for("login_page") + "?cleanup=1")
    return redirect(url_for("login_page"))


# ──────────────────────────── Main Page ────────────────────────────



# ──────────────────────────── Routes ────────────────────────────

@app.route("/health")
def health():
    return "ok", 200

@app.route("/")
@login_required
def index():
    import json
    status = hub.status()
    personas = hub.list_personas()
    persona_groups = hub.list_persona_groups()

    username = session.get("username", "")
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])

    return render_template(
        "index.html",
        ai_status=json.dumps(status),
        personas=json.dumps(personas, ensure_ascii=False),
        personaGroups=json.dumps(persona_groups, ensure_ascii=False),
        tierLimits=json.dumps(tier_limits),
        usertier=user_tier,
        username=username
    )


# ── User Persona API (simplified — no groups) ─────────────────
@app.route("/api/user/personas", methods=["GET"])
@login_required
def get_user_personas():
    """Load user's persona list from DB"""
    import json as _json
    username = session.get("username", "")
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    try:
        res = supabase_admin.table("user_personas").select("*").eq("username", username).execute()
        if res.data and len(res.data) > 0:
            config = res.data[0]
            return jsonify({
                "personas": _json.loads(config.get("persona_keys", "[]")),
                "limits": tier_limits
            })
    except Exception:
        pass
    return jsonify({"personas": [], "limits": tier_limits})


@app.route("/api/user/personas", methods=["POST"])
@login_required
def save_user_personas():
    """Save user's persona list to DB"""
    import json as _json
    username = session.get("username", "")
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    data = request.json
    persona_list = data.get("personas", [])
    if len(persona_list) > tier_limits["personas"]:
        return jsonify({"error": f"최대 {tier_limits['personas']}개 페르소나만 허용됩니다."}), 400
    try:
        payload = {
            "username": username,
            "persona_keys": _json.dumps(persona_list, ensure_ascii=False),
        }
        existing = supabase_admin.table("user_personas").select("id").eq("username", username).execute()
        if existing.data and len(existing.data) > 0:
            supabase_admin.table("user_personas").update(payload).eq("username", username).execute()
        else:
            supabase_admin.table("user_personas").insert(payload).execute()
        return jsonify({"status": "ok"})
    except Exception as e:
        print(f"  ⚠️ save_user_personas failed: {e}")
        return jsonify({"status": "ok", "fallback": True})


@app.route("/api/user/personas/create", methods=["POST"])
@login_required
def create_user_persona():
    """AI generates persona traits from name/job. Saves to user's persona list."""
    import json as _json
    username = session.get("username", "")
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    data = request.json
    name = data.get("name", "").strip()
    extra_traits = data.get("traits", "").strip()
    if not name:
        return jsonify({"error": "페르소나 이름을 입력해주세요."}), 400

    # Check current count
    current_personas = []
    try:
        res = supabase_admin.table("user_personas").select("persona_keys").eq("username", username).execute()
        if res.data and len(res.data) > 0:
            current_personas = _json.loads(res.data[0].get("persona_keys", "[]"))
    except Exception:
        local_data = data.get("current_personas", [])
        current_personas = local_data

    if len(current_personas) >= tier_limits["personas"]:
        return jsonify({"error": f"최대 {tier_limits['personas']}개 페르소나까지 생성 가능합니다."}), 400

    # AI generates persona traits from name
    trait_prompt = (
        f"You are creating a persona profile. The persona name/role is: \"{name}\".\n"
        f"{'Additional user-specified traits: ' + extra_traits if extra_traits else ''}\n\n"
        f"Generate a detailed persona profile in JSON format with these fields:\n"
        f"- prompt: A system instruction (2-3 sentences) defining this persona's expertise, personality, and communication style\n"
        f"- traits: Array of 5-8 core trait keywords (e.g. [\"analytical\",\"detail-oriented\"])\n"
        f"- skills: Array of 3-5 key skills\n"
        f"- style: Communication style in one sentence\n\n"
        f"Output ONLY valid JSON, no markdown, no explanation."
    )

    try:
        ai_response = hub.ask(trait_prompt, provider=data.get("provider", "chatgpt"))
        # AIResponse is a dataclass with .content attribute
        if not ai_response.success:
            return jsonify({"error": f"AI 응답 실패: {ai_response.error}"}), 500
        ai_text = ai_response.content.strip()
        ai_text = ai_text.strip()
        if ai_text.startswith("```"):
            ai_text = ai_text.split("\n", 1)[1] if "\n" in ai_text else ai_text[3:]
            if ai_text.endswith("```"):
                ai_text = ai_text[:-3]
            ai_text = ai_text.strip()

        try:
            profile = _json.loads(ai_text)
        except _json.JSONDecodeError:
            profile = {
                "prompt": f"You are {name}. Respond with the expertise and perspective of this role.",
                "traits": [name],
                "skills": [],
                "style": "Professional"
            }

        new_persona = {
            "key": "p_" + str(int(__import__('time').time() * 1000)),
            "name": name,
            "prompt": profile.get("prompt", f"You are {name}."),
            "traits": profile.get("traits", []),
            "skills": profile.get("skills", []),
            "style": profile.get("style", ""),
            "extra_traits": extra_traits,
        }

        current_personas.append(new_persona)

        # Save to DB
        try:
            payload = {
                "username": username,
                "persona_keys": _json.dumps(current_personas, ensure_ascii=False),
            }
            existing = supabase_admin.table("user_personas").select("id").eq("username", username).execute()
            if existing.data and len(existing.data) > 0:
                supabase_admin.table("user_personas").update(payload).eq("username", username).execute()
            else:
                supabase_admin.table("user_personas").insert(payload).execute()
        except Exception as e:
            print(f"  ⚠️ persona create DB save failed: {e}")

        return jsonify({"status": "ok", "persona": new_persona, "personas": current_personas})

    except Exception as e:
        return jsonify({"error": f"AI 페르소나 생성 실패: {str(e)}"}), 500


def _ensure_user_persona_registered(persona_key, username):
    """If persona is user-created (p_xxxx), register it in hub.PERSONAS from DB."""
    import json as _json
    if persona_key.startswith("p_") and persona_key not in hub.PERSONAS:
        try:
            res = supabase_admin.table("user_personas").select("persona_keys").eq("username", username).execute()
            if res.data and len(res.data) > 0:
                user_personas = _json.loads(res.data[0].get("persona_keys", "[]"))
                for p in user_personas:
                    if p.get("key") == persona_key:
                        hub.add_persona(persona_key, p.get("name", persona_key), p.get("prompt", f"You are {p.get('name', 'an assistant')}."))
                        return True
        except Exception:
            pass
    return persona_key in hub.PERSONAS


@app.route("/api/ask", methods=["POST"])
@login_required
def api_ask():
    try:
        data = request.json
        prompt = data.get("prompt", "")
        provider = data.get("provider", "chatgpt")
        persona = data.get("persona", "")
        response = None
        if persona:
            # Register user-created personas dynamically
            _ensure_user_persona_registered(persona, session.get("username", ""))
            memory_context = ""
            conversation_context = ""
            user_id = session.get("username", "admin")
            if supabase_client:
                try:
                    mem_result = supabase_client.table("persona_memory").select("content").eq(
                        "user_id", user_id
                    ).eq("persona_key", persona).order("created_at", desc=True).limit(20).execute()
                    if mem_result.data:
                        memories = [m["content"] for m in mem_result.data]
                        memory_context = "\n".join(f"- {m}" for m in memories)
                except Exception:
                    pass
                try:
                    conv_result = supabase_client.table("persona_conversations").select(
                        "question,answer"
                    ).eq("user_id", user_id).eq(
                        "persona_key", persona
                    ).order("created_at", desc=True).limit(10).execute()
                    if conv_result.data:
                        convs = []
                        for c in reversed(conv_result.data):
                            convs.append(f"Q: {c['question'][:200]}\nA: {c['answer'][:300]}")
                        conversation_context = "\n\n".join(convs)
                except Exception:
                    pass
            full_memory = ""
            if memory_context:
                full_memory += f"KEY INSIGHTS:\n{memory_context}"
            if conversation_context:
                if full_memory:
                    full_memory += "\n\n"
                full_memory += f"RECENT CONVERSATION HISTORY:\n{conversation_context}"
            response = hub.ask_as(prompt, persona=persona, provider=provider,
                                  memory_context=full_memory)
            if response.success and supabase_client:
                try:
                    supabase_client.table("persona_conversations").insert({
                        "user_id": user_id,
                        "persona_key": persona,
                        "question": prompt[:2000],
                        "answer": response.content[:3000],
                        "provider": provider
                    }).execute()
                except Exception:
                    pass
                if len(response.content) > 50:
                    try:
                        extract = hub.ask(
                            f"From this conversation, extract ONE concise key insight or fact worth remembering "
                            f"for future reference (1 sentence max, in the language of the content). "
                            f"If nothing worth remembering, respond with exactly 'NONE'.\n\n"
                            f"User asked: {prompt[:500]}\nResponse: {response.content[:1000]}",
                            provider=provider,
                            system_prompt="You are a memory extraction assistant. Extract only genuinely useful facts or insights. Respond with a single sentence or 'NONE'."
                        )
                        if extract.success and extract.content.strip().upper() != "NONE" and len(extract.content.strip()) > 5:
                            supabase_client.table("persona_memory").insert({
                                "user_id": user_id,
                                "persona_key": persona,
                                "content": extract.content.strip()[:500]
                            }).execute()
                    except Exception:
                        pass
        else:
            response = hub.ask(prompt, provider=provider)
        return jsonify({"provider": response.provider, "model": response.model,
                         "content": response.content, "success": response.success,
                         "error": response.error, "elapsed_seconds": response.elapsed_seconds})
    except Exception as e:
        return jsonify({"success": False, "error": f"서버 오류: {str(e)}", "content": "", "provider": "error", "model": "", "elapsed_seconds": 0}), 500

@app.route("/api/ask_stream", methods=["POST"])
@login_required
def api_ask_stream():
    data = request.json
    prompt = data.get("prompt", "")
    provider = data.get("provider", "chatgpt")
    persona = data.get("persona", "")
    chat_context = data.get("chat_context", [])
    user_id = session.get("username", "admin")

    # Build conversation history context from frontend messages
    history_context = ""
    if chat_context and isinstance(chat_context, list):
        history_lines = []
        for msg in chat_context[-10:]:
            role = msg.get("role", "user")
            speaker = msg.get("speaker", "")
            content = msg.get("content", "")[:500]
            if role == "user":
                history_lines.append(f"User: {content}")
            else:
                history_lines.append(f"{speaker}: {content}")
        if history_lines:
            history_context = "PREVIOUS CONVERSATION (for context continuity):\n" + "\n".join(history_lines) + "\n\nNow respond to the latest message:"

    # Register user-created personas dynamically
    if persona:
        _ensure_user_persona_registered(persona, session.get("username", ""))

    def generate():
        try:
            full_content = ""
            if persona:
                memory_context = ""
                conversation_context = ""
                if supabase_client:
                    try:
                        mem_result = supabase_client.table("persona_memory").select("content").eq(
                            "user_id", user_id
                        ).eq("persona_key", persona).order("created_at", desc=True).limit(20).execute()
                        if mem_result.data:
                            memories = [m["content"] for m in mem_result.data]
                            memory_context = "\n".join(f"- {m}" for m in memories)
                    except Exception:
                        pass
                    try:
                        conv_result = supabase_client.table("persona_conversations").select(
                            "question,answer"
                        ).eq("user_id", user_id).eq(
                            "persona_key", persona
                        ).order("created_at", desc=True).limit(10).execute()
                        if conv_result.data:
                            convs = []
                            for c in reversed(conv_result.data):
                                convs.append(f"Q: {c['question'][:200]}\nA: {c['answer'][:300]}")
                            conversation_context = "\n\n".join(convs)
                    except Exception:
                        pass
                full_memory = ""
                if memory_context:
                    full_memory += f"KEY INSIGHTS:\n{memory_context}"
                if conversation_context:
                    if full_memory:
                        full_memory += "\n\n"
                    full_memory += f"RECENT CONVERSATION HISTORY:\n{conversation_context}"
                if history_context:
                    if full_memory:
                        full_memory += "\n\n"
                    full_memory += history_context
                
                for chunk in hub.ask_as_stream(prompt, persona=persona, provider=provider, memory_context=full_memory):
                    if chunk:
                        full_content += chunk
                        yield chunk

                if supabase_client:
                    try:
                        supabase_client.table("persona_conversations").insert({
                            "user_id": user_id,
                            "persona_key": persona,
                            "question": prompt[:2000],
                            "answer": full_content[:3000],
                            "provider": provider
                        }).execute()
                    except Exception:
                        pass
            else:
                # Inject chat context as system prompt for cross-provider continuity
                sys_prompt = history_context if history_context else ""
                for chunk in hub.ask_stream(prompt, provider=provider, system_prompt=sys_prompt):
                    if chunk:
                        full_content += chunk
                        yield chunk
        except Exception as e:
            yield f"\n[Stream Error: {str(e)}]"

    from flask import Response, stream_with_context
    return Response(stream_with_context(generate()), mimetype='text/plain')

@app.route("/api/compare", methods=["POST"])
@login_required
def api_compare():
    try:
        data = request.json
        results = hub.ask_all(data.get("prompt", ""))
        return jsonify({"results": [{"provider": r.provider, "model": r.model,
            "content": r.content, "success": r.success, "error": r.error,
            "elapsed_seconds": r.elapsed_seconds} for r in results]})
    except Exception as e:
        return jsonify({"error": str(e), "results": []}), 500


@app.route("/api/debate", methods=["POST"])
@login_required
def api_debate():
    try:
        topic = request.json.get("topic", "")
        av = hub.available_providers()
        def generate():
            try:
                result = hub.debate(topic=topic, rounds=2, ai_for=av[0],
                    ai_against=av[1] if len(av) > 1 else av[0],
                    judge=av[2] if len(av) >= 3 else av[0])
                yield f"data: {json.dumps(result, ensure_ascii=False, default=str)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'error': str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/api/discuss", methods=["POST"])
@login_required
def api_discuss():
    try:
        topic = request.json.get("topic", "")
        def generate():
            try:
                for chunk in hub.discuss_stream(topic=topic, rounds=2):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/best", methods=["POST"])
@login_required
def api_best():
    try:
        question = request.json.get("question", "")
        def generate():
            try:
                for chunk in hub.find_best_stream(question=question):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona_debate", methods=["POST"])
@login_required
def api_persona_debate():
    try:
        data = request.json
        username = session.get("user", "")
        pf = data.get("persona_for", "elon_musk")
        pa = data.get("persona_against", "trump")
        _ensure_user_persona_registered(pf, username)
        _ensure_user_persona_registered(pa, username)
        def generate():
            try:
                for chunk in hub.persona_debate_stream(
                    topic=data.get("topic", ""),
                    persona_for=pf,
                    persona_against=pa,
                    rounds=2
                ):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona_discuss", methods=["POST"])
@login_required
def api_persona_discuss():
    try:
        data = request.json
        personas = data.get("personas", [])
        if len(personas) < 2:
            return jsonify({"error": "Select at least 2 personas for discussion"}), 400
        username = session.get("user", "")
        for pk in personas:
            _ensure_user_persona_registered(pk, username)
        topic = data.get("topic", "")
        def generate():
            try:
                for chunk in hub.persona_discuss_stream(topic=topic, persona_keys=personas, rounds=2):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona_report", methods=["POST"])
@login_required
def api_persona_report():
    try:
        data = request.json
        personas = data.get("personas", [])
        if len(personas) < 1:
            return jsonify({"error": "Select at least 1 persona"}), 400
        topic = data.get("topic", "")
        provider = data.get("provider", "chatgpt")
        def generate():
            try:
                for chunk in hub.multi_persona_report_stream(topic=topic, persona_keys=personas, provider=provider):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona_chain", methods=["POST"])
@login_required
def api_persona_chain():
    try:
        data = request.json
        personas = data.get("personas", [])
        if len(personas) < 2:
            return jsonify({"error": "Select at least 2 personas for chain"}), 400
        topic = data.get("topic", "")
        provider = data.get("provider", "chatgpt")
        def generate():
            try:
                for chunk in hub.persona_chain_stream(topic=topic, persona_keys=personas, provider=provider):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona_vote", methods=["POST"])
@login_required
def api_persona_vote():
    try:
        data = request.json
        personas = data.get("personas", [])
        if len(personas) < 2:
            return jsonify({"error": "Select at least 2 personas for voting"}), 400
        proposal = data.get("proposal", data.get("topic", ""))
        provider = data.get("provider", "chatgpt")
        def generate():
            try:
                for chunk in hub.persona_vote_stream(proposal=proposal, persona_keys=personas, provider=provider):
                    yield f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n"
            except Exception as e:
                yield f"data: {json.dumps({'type':'error','error':str(e)})}\n\n"
            yield "data: [DONE]\n\n"
        return Response(stream_with_context(generate()), content_type="text/event-stream",
                        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/decision_matrix", methods=["POST"])
@login_required
def api_decision_matrix():
    """Evaluate options against criteria using selected personas."""
    data = request.json
    options = data.get("options", [])
    criteria = data.get("criteria", [])
    personas = data.get("personas", [])
    if len(options) < 2 or len(criteria) < 1 or len(personas) < 1:
        return jsonify({"error": "Need ≥2 options, ≥1 criteria, ≥1 persona"}), 400

    available = [p for p in ["chatgpt", "gemini", "azure", "claude", "grok"]
                 if p in hub.available_providers()]
    if not available:
        return jsonify({"error": "No AI providers available"}), 400

    scores = []
    for i, pkey in enumerate(personas):
        name = hub.get_persona_name(pkey)
        prompt_text = hub.get_persona_prompt(pkey)
        if not prompt_text:
            continue
        ai = available[i % len(available)]
        options_str = ", ".join(options)
        criteria_str = ", ".join(criteria)
        sys_prompt = (
            f"{prompt_text}\n\n"
            f"You are evaluating options in a decision matrix. "
            f"Score each option on each criterion from 1-10.\n"
            f"Respond in EXACTLY this format for each option:\n"
            f"OPTION_NAME: criterion1=N, criterion2=N, ...\n"
            f"Then one line: BEST: [your recommended option]"
        )
        resp = hub.ask(
            f"Evaluate these options: {options_str}\n"
            f"Against these criteria: {criteria_str}\n"
            f"Score each 1-10.",
            provider=ai, system_prompt=sys_prompt
        )
        scores.append({
            "persona_key": pkey, "persona_name": name,
            "evaluation": resp.content if resp.success else "Error",
            "provider": ai
        })

    # Synthesize
    all_evals = "\n\n".join(f"[{s['persona_name']}]:\n{s['evaluation']}" for s in scores)
    synth = hub.ask(
        f"Synthesize these decision matrix evaluations:\n\n{all_evals}\n\n"
        f"Options: {', '.join(options)}\nCriteria: {', '.join(criteria)}\n\n"
        f"Create a final scorecard table and recommend the best option.",
        provider=available[0],
        system_prompt="You are a decision matrix synthesizer. Present a clear scorecard."
    )
    return jsonify({
        "options": options, "criteria": criteria,
        "evaluations": scores, "synthesis": synth.content if synth.success else synth.error
    })


# ── Prompt Library ──

@app.route("/api/prompts", methods=["GET"])
@login_required
def api_prompts_list():
    if not supabase_client:
        return jsonify({"prompts": []})
    try:
        q = request.args.get("q", "").strip()
        category = request.args.get("category", "").strip()
        uid = session.get("username", "admin")
        query = supabase_client.table("saved_prompts").select("*").eq("user_id", uid)
        if category:
            query = query.eq("category", category)
        result = query.order("created_at", desc=True).execute()
        prompts = result.data or []
        if q:
            ql = q.lower()
            prompts = [p for p in prompts if ql in p.get("name","").lower() or ql in p.get("prompt","").lower()]
        return jsonify({"prompts": prompts})
    except Exception:
        return jsonify({"prompts": []})


@app.route("/api/prompts/categories", methods=["GET"])
@login_required
def api_prompts_categories():
    if not supabase_client:
        return jsonify({"categories": []})
    try:
        uid = session.get("username", "admin")
        result = supabase_client.table("saved_prompts").select("category").eq("user_id", uid).execute()
        cats = list({r.get("category") for r in (result.data or []) if r.get("category")})
        return jsonify({"categories": sorted(cats)})
    except Exception:
        return jsonify({"categories": []})


@app.route("/api/prompts", methods=["POST"])
@login_required
def api_prompts_save():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json
    row = {
        "user_id": session.get("username", "admin"),
        "name": data.get("name", "Untitled"),
        "prompt": data.get("prompt", ""),
        "mode": data.get("mode", "chat"),
        "personas": data.get("personas", []),
    }
    if data.get("category"):
        row["category"] = data["category"]
    try:
        result = supabase_client.table("saved_prompts").insert(row).execute()
    except Exception:
        row.pop("category", None)
        result = supabase_client.table("saved_prompts").insert(row).execute()
    return jsonify({"prompt": result.data[0] if result.data else {}, "success": True})


@app.route("/api/prompts/<prompt_id>", methods=["PATCH"])
@login_required
def api_prompts_update(prompt_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json
    allowed = {k: v for k, v in data.items() if k in ("name", "category", "use_count")}
    if not allowed:
        return jsonify({"error": "Nothing to update"}), 400
    try:
        supabase_client.table("saved_prompts").update(allowed).eq("id", prompt_id).execute()
    except Exception:
        allowed.pop("category", None)
        if allowed:
            supabase_client.table("saved_prompts").update(allowed).eq("id", prompt_id).execute()
    return jsonify({"updated": True})


@app.route("/api/prompts/<prompt_id>", methods=["DELETE"])
@login_required
def api_prompts_delete(prompt_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    supabase_client.table("saved_prompts").delete().eq("id", prompt_id).execute()
    return jsonify({"deleted": True})


@app.route("/api/add_custom_persona", methods=["POST"])
@login_required
def api_add_custom_persona():
    data = request.json
    key = data.get("key", "custom")
    name = data.get("name", "Custom")
    prompt_text = data.get("prompt", "You are a helpful assistant.")
    hub.add_persona(key, name, prompt_text)
    return jsonify({"success": True, "key": key, "name": name})


@app.route("/api/upload", methods=["POST"])
@login_required
def api_upload():
    if 'file' not in request.files:
        return jsonify({"success": False, "error": "No file"})
    file = request.files['file']
    if not file.filename:
        return jsonify({"success": False, "error": "No filename"})
    ext = os.path.splitext(file.filename)[1].lower()
    filepath = os.path.join(UPLOAD_DIR, file.filename)
    file.save(filepath)
    try:
        content = ""
        if ext in (".txt",".md",".csv",".json",".py",".js",".html",".css",".xml",".log",
                    ".yaml",".yml",".toml",".ini",".cfg",".bat",".sh",".sql"):
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
        elif ext == ".pdf":
            # Step 1: Try normal text extraction
            try:
                import pypdf
                with open(filepath, "rb") as f:
                    reader = pypdf.PdfReader(f)
                    for page in reader.pages: content += page.extract_text() or ""
            except Exception: content = ""

            # Step 2: If text is empty → AI Vision OCR using GPT-4o
            if len(content.strip()) < 50 and hub.openai_api_key:
                try:
                    import fitz, base64  # PyMuPDF
                    from openai import OpenAI
                    client = OpenAI(api_key=hub.openai_api_key)
                    doc = fitz.open(filepath)
                    ocr_pages = []
                    total_pages = len(doc)
                    for i in range(total_pages):
                        pix = doc[i].get_pixmap(dpi=150)
                        b64 = base64.b64encode(pix.tobytes("png")).decode()
                        resp = client.chat.completions.create(
                            model="gpt-4o",
                            messages=[{"role":"user","content":[
                                {"type":"text","text":f"이 PDF 페이지({i+1}/{total_pages})의 모든 텍스트를 정확하게 추출해주세요. 표, 숫자, 특수문자도 포함해서 원본 그대로 출력해주세요."},
                                {"type":"image_url","image_url":{"url":f"data:image/png;base64,{b64}","detail":"high"}}
                            ]}],
                            max_tokens=4096
                        )
                        ocr_pages.append(f"=== 페이지 {i+1} ===\n{resp.choices[0].message.content}")
                    content = "\n\n".join(ocr_pages)
                    doc.close()
                except Exception as e:
                    if not content.strip(): content = f"[OCR 실패: {str(e)}]"
            elif len(content.strip()) < 50:
                content = "[PDF: 이미지 기반 PDF - OpenAI API 키 필요]"
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(filepath)
                content = "\n".join(p.text for p in doc.paragraphs)
            except ImportError: content = "[DOCX: pip install python-docx]"
        elif ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(filepath, read_only=True)
                for s in wb.sheetnames:
                    ws = wb[s]; content += f"\n--- {s} ---\n"
                    for row in ws.iter_rows(values_only=True):
                        content += "\t".join(str(c) if c else "" for c in row) + "\n"
            except ImportError: content = "[XLSX: pip install openpyxl]"
        else:
            try:
                with open(filepath, "r", encoding="utf-8", errors="replace") as f: content = f.read()
            except: content = "[Unsupported format]"
        chunks_indexed = 0
        if content.strip():
            try:
                chunks_indexed = hub.index_document(text=content, document_id=file.filename)
            except Exception as e:
                print(f"Warning: Failed to index document for RAG: {e}")

        if len(content) > 100000: content = content[:100000] + "\n\n[... truncated ...]"
        return jsonify({
            "success": True, 
            "filename": file.filename,
            "size": os.path.getsize(filepath), 
            "char_count": len(content), 
            "content": content,
            "chunks_indexed": chunks_indexed
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})

@app.route("/api/query_rag", methods=["POST"])
@login_required
def api_query_rag():
    """Query ChromaDB for relevant exact chunks for uploaded files before sending to AI."""
    try:
        data = request.json
        query = data.get("query", "")
        files = data.get("files", [])
        top_k = data.get("top_k", 3)
        
        if not query or not files:
            return jsonify({"success": False, "context": ""})
            
        combined_context = []
        for f in files:
            snippet = hub.query_document(query, document_id=f, top_k=top_k)
            if snippet:
                combined_context.append(f"=== 관련 문서 일부: {f} ===\n{snippet}")
                
        # If no RAG results were found (e.g. indexing failed), fallback to empty context
        final_context = "\n\n".join(combined_context) if combined_context else ""
        return jsonify({"success": True, "context": final_context})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})



@app.route("/api/visualize", methods=["POST"])
@login_required
def api_visualize():
    """Ask AI to analyze context and return Chart.js data or Mermaid diagram"""
    try:
        context = request.json.get("context", "").strip()
        if not context:
            return jsonify({"success": False, "error": "No context provided"})

        prompt = f"""다음 데이터/내용을 분석하고, 가장 적합한 시각화를 선택하세요.

데이터:
{context[:6000]}

다음 중 하나의 JSON 형식으로만 응답하세요 (다른 텍스트 없이):

1. 숫자 데이터 → Chart.js 차트:
{{"type": "chart", "title": "제목", "chart_data": {{"type": "bar"|"line"|"pie"|"doughnut", "labels": [...], "datasets": [{{"label": "...", "data": [...]}}]}}}}

2. 프로세스/관계 → Mermaid 다이어그램:
{{"type": "diagram", "title": "제목", "diagram": "flowchart TD\\n  A --> B"}}

JSON만 반환하세요."""

        from openai import OpenAI
        client = OpenAI(api_key=hub.openai_api_key)
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500,
            response_format={"type": "json_object"}
        )
        import json as _json
        result = _json.loads(resp.choices[0].message.content)

        if result.get("type") == "chart":
            return jsonify({"success": True, "title": result.get("title", "Chart"),
                            "chart_data": result["chart_data"]})
        elif result.get("type") == "diagram":
            return jsonify({"success": True, "title": result.get("title", "Diagram"),
                            "diagram": result["diagram"]})
        else:
            return jsonify({"success": False, "error": "예상치 못한 응답 형식"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


@app.route("/api/fetch_url", methods=["POST"])
@login_required
def api_fetch_url():
    """Fetch and extract text content from a URL for AI analysis"""
    try:
        url = request.json.get("url", "").strip()
        if not url:
            return jsonify({"success": False, "error": "No URL provided"})
        if not url.startswith(("http://", "https://")):
            url = "https://" + url

        import requests as req
        from bs4 import BeautifulSoup

        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}
        resp = req.get(url, headers=headers, timeout=15)
        resp.raise_for_status()

        soup = BeautifulSoup(resp.text, "html.parser")
        # Remove scripts, styles, nav elements
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "iframe"]):
            tag.decompose()

        text = soup.get_text(separator="\n", strip=True)
        # Clean up excessive blank lines
        import re
        text = re.sub(r'\n{3,}', '\n\n', text)

        if len(text) > 40000:
            text = text[:40000] + "\n\n[... 내용이 길어 일부 생략됨 ...]"

        title = soup.title.string.strip() if soup.title else url

        return jsonify({
            "success": True,
            "url": url,
            "title": title,
            "content": text,
            "char_count": len(text)
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)})


# ──────────────────────────── Conversation History (Supabase) ────────────────────────────

@app.route("/api/conversations", methods=["GET"])
@login_required
def api_conversations_list():
    """List all conversations for current user"""
    if not supabase_client:
        return jsonify({"conversations": [], "supabase": False})
    try:
        username = session.get("username", APP_USERNAME)
        result = supabase_client.table("conversations").select("*").eq(
            "username", username).order("updated_at", desc=True).limit(50).execute()
        return jsonify({"conversations": result.data, "supabase": True})
    except Exception as e:
        return jsonify({"conversations": [], "error": str(e), "supabase": True})


@app.route("/api/conversations", methods=["POST"])
@login_required
def api_conversations_create():
    """Create a new conversation"""
    if not supabase_client:
        return jsonify({"error": "Supabase not configured"}), 400
    try:
        data = request.json
        username = session.get("username", APP_USERNAME)
        result = supabase_client.table("conversations").insert({
            "title": data.get("title", "New Chat"),
            "mode": data.get("mode", "chat"),
            "username": username,
        }).execute()
        return jsonify({"conversation": result.data[0]})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/conversations/<conv_id>/messages", methods=["GET"])
@login_required
def api_conversation_messages(conv_id):
    """Get all messages for a conversation"""
    if not supabase_client:
        return jsonify({"messages": []})
    try:
        result = supabase_client.table("messages").select("*").eq(
            "conversation_id", conv_id).order("created_at").execute()
        return jsonify({"messages": result.data})
    except Exception as e:
        return jsonify({"messages": [], "error": str(e)})


@app.route("/api/conversations/<conv_id>/messages", methods=["POST"])
@login_required
def api_conversation_save_message(conv_id):
    """Save a message to a conversation"""
    if not supabase_client:
        return jsonify({"saved": False})
    try:
        data = request.json
        msg = {
            "conversation_id": conv_id,
            "role": data.get("role", "user"),
            "speaker": data.get("speaker", ""),
            "content": data.get("content", ""),
            "provider": data.get("provider", ""),
            "model": data.get("model", ""),
            "badge": data.get("badge", ""),
            "elapsed_seconds": data.get("elapsed_seconds"),
        }
        supabase_client.table("messages").insert(msg).execute()
        # Update conversation title & timestamp
        title_update = {"updated_at": datetime.utcnow().isoformat()}
        if data.get("role") == "user" and data.get("update_title"):
            title_text = data.get("content", "")[:50]
            title_update["title"] = title_text
        supabase_client.table("conversations").update(title_update).eq("id", conv_id).execute()
        return jsonify({"saved": True})
    except Exception as e:
        return jsonify({"saved": False, "error": str(e)})


@app.route("/api/conversations/<conv_id>", methods=["DELETE"])
@login_required
def api_conversation_delete(conv_id):
    """Delete a conversation and all its messages"""
    if not supabase_client:
        return jsonify({"deleted": False})
    try:
        supabase_client.table("conversations").delete().eq("id", conv_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"deleted": False, "error": str(e)})


# ──────────────────────────── Workspace: Folders & Files ────────────────────────────

@app.route("/api/folders", methods=["GET"])
@login_required
def api_folders_list():
    """List all folders with workspace usage info"""
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    if not supabase_client:
        return jsonify({"folders": [], "workspace": False, "limits": tier_limits})
    try:
        username = session.get("username", "admin")
        result = supabase_client.table("folders").select("*").eq(
            "user_id", username
        ).order("created_at", desc=False).execute()
        # Calculate usage (lightweight — id only)
        folders = result.data or []
        file_result = supabase_client.table("workspace_files").select("id").eq("user_id", username).execute()
        files = file_result.data or []
        usage = {
            "folders": len(folders),
            "files": len(files),
            "storage_mb": 0  # calculated on-demand during file creation
        }
        return jsonify({"folders": folders, "workspace": True, "limits": tier_limits, "usage": usage})
    except:
        return jsonify({"folders": [], "workspace": True, "limits": tier_limits})


@app.route("/api/folders", methods=["POST"])
@login_required
def api_folders_create():
    """Create a new folder (tier-limited)"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    # Free tier: no workspace
    if tier_limits["folders"] == 0:
        return jsonify({"error": "워크스페이스는 Premium 이상에서 사용 가능합니다. 업그레이드해주세요."}), 403
    try:
        # Check current folder count
        username = session.get("username", "admin")
        existing = supabase_client.table("folders").select("id").eq("user_id", username).execute()
        current_count = len(existing.data) if existing.data else 0
        if current_count >= tier_limits["folders"]:
            return jsonify({"error": f"폴더 최대 {tier_limits['folders']}개까지 생성 가능합니다. ({current_count}/{tier_limits['folders']})"}), 400
        data = request.json
        result = supabase_client.table("folders").insert({
            "user_id": username,
            "name": data.get("name", "New Folder"),
            "icon": data.get("icon", "📁"),
            "description": data.get("description", "")
        }).execute()
        return jsonify({"folder": result.data[0] if result.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/<folder_id>", methods=["PUT"])
@login_required
def api_folders_update(folder_id):
    """Rename/update folder"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        data = request.json
        update = {}
        if "name" in data: update["name"] = data["name"]
        if "icon" in data: update["icon"] = data["icon"]
        if "description" in data: update["description"] = data["description"]
        supabase_client.table("folders").update(update).eq("id", folder_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/<folder_id>", methods=["DELETE"])
@login_required
def api_folders_delete(folder_id):
    """Delete folder and its files"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("workspace_files").delete().eq("folder_id", folder_id).execute()
        supabase_client.table("folders").delete().eq("id", folder_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/<folder_id>/files", methods=["GET"])
@login_required
def api_files_list(folder_id):
    """List files in a folder"""
    if not supabase_client:
        return jsonify({"files": []})
    try:
        result = supabase_client.table("workspace_files").select("*").eq(
            "folder_id", folder_id
        ).order("updated_at", desc=True).execute()
        return jsonify({"files": result.data})
    except:
        return jsonify({"files": []})


@app.route("/api/folders/<folder_id>/files", methods=["POST"])
@login_required
def api_files_create(folder_id):
    """Create/save a file in a folder (tier-limited)"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    user_tier = session.get("user_tier", "free")
    tier_limits = hub.TIER_LIMITS.get(user_tier, hub.TIER_LIMITS["free"])
    # Free tier: no workspace
    if tier_limits["files"] == 0:
        return jsonify({"error": "파일 저장은 Premium 이상에서 사용 가능합니다."}), 403
    try:
        data = request.json
        username = session.get("username", "admin")
        content = data.get("content", "")
        # Check file size limit
        content_size_mb = len(content.encode("utf-8")) / (1024 * 1024)
        if content_size_mb > tier_limits["file_size_mb"]:
            return jsonify({"error": f"파일 크기 제한 초과 ({tier_limits['file_size_mb']}MB 이하만 허용)"}), 400
        # Check total file count (lightweight — id only)
        all_files = supabase_client.table("workspace_files").select("id").eq("user_id", username).execute()
        current_file_count = len(all_files.data) if all_files.data else 0
        if current_file_count >= tier_limits["files"]:
            return jsonify({"error": f"파일 최대 {tier_limits['files']}개까지 저장 가능합니다. ({current_file_count}/{tier_limits['files']})"}), 400
        # Check total storage (only when creating, fetch content for size calc)
        storage_result = supabase_client.table("workspace_files").select("content").eq("user_id", username).execute()
        total_bytes = sum(len((f.get("content") or "").encode("utf-8")) for f in (storage_result.data or []))
        total_mb = total_bytes / (1024 * 1024)
        if total_mb + content_size_mb > tier_limits["storage_mb"]:
            return jsonify({"error": f"저장 공간 한도 초과 ({total_mb:.1f}/{tier_limits['storage_mb']}MB 사용 중)"}), 400
        result = supabase_client.table("workspace_files").insert({
            "folder_id": folder_id,
            "user_id": username,
            "name": data.get("name", "Untitled"),
            "type": data.get("type", "note"),  # note, conversation, slides, file
            "content": content,
            "metadata": data.get("metadata", {})
        }).execute()
        return jsonify({"file": result.data[0] if result.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/workspace/search", methods=["GET"])
@login_required
def api_workspace_search():
    """Search files across all folders by name or tag"""
    if not supabase_client:
        return jsonify({"files": []})
    try:
        username = session.get("username", "admin")
        q = request.args.get("q", "").strip().lower()
        tag = request.args.get("tag", "").strip().lower()
        result = supabase_client.table("workspace_files").select(
            "id,name,type,folder_id,created_at,updated_at,metadata"
        ).eq("user_id", username).execute()
        files = result.data or []
        if q:
            files = [f for f in files if q in f.get("name", "").lower()]
        if tag:
            files = [f for f in files if tag in [
                t.lower() for t in ((f.get("metadata") or {}).get("tags") or [])
            ]]
        # Attach folder info
        folder_ids = list({f["folder_id"] for f in files if f.get("folder_id")})
        folder_map = {}
        if folder_ids:
            fr = supabase_client.table("folders").select("id,name,icon").in_("id", folder_ids).execute()
            for fol in (fr.data or []):
                folder_map[fol["id"]] = fol
        for f in files:
            fol = folder_map.get(f.get("folder_id"), {})
            f["folder_name"] = fol.get("name", "")
            f["folder_icon"] = fol.get("icon", "📁")
        files.sort(key=lambda f: f.get("updated_at") or "", reverse=True)
        return jsonify({"files": files, "count": len(files)})
    except Exception as e:
        return jsonify({"files": [], "error": str(e)})


@app.route("/api/files/<file_id>", methods=["GET"])
@login_required
def api_files_get(file_id):
    """Get file content"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        result = supabase_client.table("workspace_files").select("*").eq("id", file_id).execute()
        if result.data:
            return jsonify({"file": result.data[0]})
        return jsonify({"error": "Not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/files/<file_id>", methods=["PUT"])
@login_required
def api_files_update(file_id):
    """Update file content"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        data = request.json
        update = {"updated_at": "now()"}
        if "name" in data: update["name"] = data["name"]
        if "content" in data: update["content"] = data["content"]
        if "metadata" in data: update["metadata"] = data["metadata"]
        if "tags" in data:
            # Merge tags into existing metadata
            existing = supabase_client.table("workspace_files").select("metadata").eq("id", file_id).execute()
            meta = (existing.data[0].get("metadata") or {}) if existing.data else {}
            meta["tags"] = [t.strip() for t in data["tags"] if t.strip()]
            update["metadata"] = meta
        supabase_client.table("workspace_files").update(update).eq("id", file_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/files/<file_id>", methods=["DELETE"])
@login_required
def api_files_delete(file_id):
    """Delete a file"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("workspace_files").delete().eq("id", file_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ──────────────────────────── Persona Memory API ────────────────────────────


@app.route("/api/persona/<persona_key>/memory", methods=["GET"])
@login_required
def api_persona_memory_list(persona_key):
    """List memories for a persona"""
    if not supabase_client:
        return jsonify({"memories": [], "error": "No database"}), 200
    try:
        result = supabase_client.table("persona_memory").select("*").eq(
            "user_id", session.get("username", "admin")
        ).eq("persona_key", persona_key).order("created_at", desc=True).execute()
        return jsonify({"memories": result.data or [], "count": len(result.data or [])})
    except Exception as e:
        return jsonify({"memories": [], "error": str(e)}), 200


@app.route("/api/persona/<persona_key>/memory", methods=["POST"])
@login_required
def api_persona_memory_add(persona_key):
    """Manually add a memory for a persona"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        data = request.json
        content = data.get("content", "").strip()
        if not content:
            return jsonify({"error": "Empty memory"}), 400
        result = supabase_client.table("persona_memory").insert({
            "user_id": session.get("username", "admin"),
            "persona_key": persona_key,
            "content": content[:500]
        }).execute()
        return jsonify({"memory": result.data[0] if result.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona/memory/<memory_id>", methods=["DELETE"])
@login_required
def api_persona_memory_delete(memory_id):
    """Delete a specific memory"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("persona_memory").delete().eq("id", memory_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona/<persona_key>/memory/clear", methods=["DELETE"])
@login_required
def api_persona_memory_clear(persona_key):
    """Clear all memories AND conversations for a persona"""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("persona_memory").delete().eq(
            "user_id", session.get("username", "admin")
        ).eq("persona_key", persona_key).execute()
        supabase_client.table("persona_conversations").delete().eq(
            "user_id", session.get("username", "admin")
        ).eq("persona_key", persona_key).execute()
        return jsonify({"cleared": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/persona/<persona_key>/conversations", methods=["GET"])
@login_required
def api_persona_conversations(persona_key):
    """List recent conversations for a persona"""
    if not supabase_client:
        return jsonify({"conversations": []})
    try:
        result = supabase_client.table("persona_conversations").select("*").eq(
            "user_id", session.get("username", "admin")
        ).eq("persona_key", persona_key).order("created_at", desc=True).limit(20).execute()
        return jsonify({"conversations": result.data or [], "count": len(result.data or [])})
    except Exception:
        return jsonify({"conversations": []})


# ──────────────────────────── Image Generation: DALL-E 3 ────────────────────────────

@app.route("/api/generate-image", methods=["POST"])
@login_required
def api_generate_image():
    """Generate images using OpenAI DALL-E 3"""
    try:
        from openai import OpenAI
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return jsonify({"error": "OpenAI API key not configured"}), 400
        client = OpenAI(api_key=api_key)
        data = request.json
        prompt = data.get("prompt", "").strip()
        if not prompt:
            return jsonify({"error": "Prompt is required"}), 400
        size = data.get("size", "1024x1024")   # 1024x1024, 1024x1792, 1792x1024
        quality = data.get("quality", "standard")  # standard, hd
        style = data.get("style", "vivid")  # vivid, natural
        # Validate size
        valid_sizes = ["1024x1024", "1024x1792", "1792x1024"]
        if size not in valid_sizes:
            size = "1024x1024"
        response = client.images.generate(
            model="dall-e-3",
            prompt=prompt,
            size=size,
            quality=quality,
            style=style,
            n=1,
        )
        image_url = response.data[0].url
        revised_prompt = response.data[0].revised_prompt
        return jsonify({
            "success": True,
            "image_url": image_url,
            "revised_prompt": revised_prompt,
            "size": size,
            "quality": quality,
            "style": style,
        })
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ──────────────────────────── Audio: OpenAI TTS ────────────────────────────

@app.route("/api/tts", methods=["POST"])
@login_required
def api_tts():
    """Convert text to natural speech using OpenAI TTS"""
    try:
        from openai import OpenAI
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return jsonify({"error": "OpenAI API key not configured"}), 400
        client = OpenAI(api_key=api_key)
        data = request.json
        text = data.get("text", "")[:4096]  # TTS limit
        voice = data.get("voice", "alloy")  # alloy, echo, fable, onyx, nova, shimmer
        response = client.audio.speech.create(
            model="tts-1",
            voice=voice,
            input=text
        )
        # Return audio as binary
        from flask import Response
        return Response(response.content, mimetype="audio/mpeg",
                       headers={"Content-Disposition": "inline"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ──────────────────────────── Audio: Whisper Transcription ────────────────────────────

@app.route("/api/transcribe", methods=["POST"])
@login_required
def api_transcribe():
    """Transcribe audio file to text using OpenAI Whisper"""
    try:
        from openai import OpenAI
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return jsonify({"error": "OpenAI API key not configured"}), 400
        client = OpenAI(api_key=api_key)

        if "file" not in request.files:
            return jsonify({"error": "No audio file uploaded"}), 400

        audio_file = request.files["file"]
        # Save to temp file
        suffix = os.path.splitext(audio_file.filename)[1] or ".mp3"
        tmp_path = os.path.join(UPLOAD_DIR, f"whisper_{secrets.token_hex(4)}{suffix}")
        audio_file.save(tmp_path)

        try:
            with open(tmp_path, "rb") as f:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=f
                )
            return jsonify({"text": transcript.text, "success": True})
        finally:
            os.remove(tmp_path)
    except Exception as e:
        return jsonify({"error": str(e), "success": False}), 500


# ──────────────────────────── Slides: PPTX Generation ────────────────────────────

@app.route("/api/slides", methods=["POST"])
@login_required
def api_slides():
    """Generate PPTX from slide data and return for download"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from io import BytesIO

        data = request.json
        slides_data = data.get("slides", [])
        title = data.get("title", "AI Hub Presentation")

        if not slides_data:
            return jsonify({"error": "No slides data"}), 400

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for i, slide_data in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Dark gradient background
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0x0a, 0x0a, 0x16)

            # Top accent bar (gradient effect via shapes)
            accent_bar = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, Inches(0.06)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(0x6c, 0x5c, 0xe7)
            accent_bar.line.fill.background()

            s_title = slide_data.get("title", "")
            s_content = slide_data.get("content", "")
            s_bullets = slide_data.get("bullets", [])

            # Title
            left = Inches(0.8)
            top = Inches(0.6) if i == 0 else Inches(0.4)
            width = Inches(11.7)
            height = Inches(1.8) if i == 0 else Inches(1.0)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s_title
            p.font.size = Pt(44) if i == 0 else Pt(30)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xf0, 0xf0, 0xff)
            p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER

            # Subtitle for title slide
            if i == 0 and s_content:
                p2 = tf.add_paragraph()
                p2.text = s_content
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(0xa2, 0x9b, 0xfe)
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(12)

            # Accent line under title (non-title slides)
            if i > 0:
                line = slide.shapes.add_shape(
                    1, Inches(0.8), Inches(1.4), Inches(2.0), Inches(0.03)
                )
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(0x6c, 0x5c, 0xe7)
                line.line.fill.background()

            # Bullets
            if s_bullets and i > 0:
                b_top = Inches(1.7)
                b_height = Inches(5.0)
                txBox2 = slide.shapes.add_textbox(left, b_top, width, b_height)
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for j, bullet in enumerate(s_bullets):
                    p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    p.text = f"▸  {bullet}"
                    p.font.size = Pt(20)
                    p.font.color.rgb = RGBColor(0xd1, 0xd5, 0xdb)
                    p.space_after = Pt(10)

            # Content paragraph for non-title slides
            if s_content and i > 0 and not s_bullets:
                c_top = Inches(1.7)
                txBox3 = slide.shapes.add_textbox(left, c_top, width, Inches(5.0))
                tf3 = txBox3.text_frame
                tf3.word_wrap = True
                p = tf3.paragraphs[0]
                p.text = s_content
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0xd1, 0xd5, 0xdb)
                p.line_spacing = Pt(28)

            # Footer: page number
            footer = slide.shapes.add_textbox(
                Inches(12.0), Inches(7.0), Inches(1.0), Inches(0.3)
            )
            ftf = footer.text_frame
            fp = ftf.paragraphs[0]
            fp.text = f"{i + 1}"
            fp.font.size = Pt(11)
            fp.font.color.rgb = RGBColor(0x6b, 0x72, 0x80)
            fp.alignment = PP_ALIGN.RIGHT

        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        from flask import send_file
        safe_title = "".join(c for c in title if c.isalnum() or c in " -_").strip()[:50] or "presentation"
        return send_file(buffer, mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        as_attachment=True, download_name=f"{safe_title}.pptx")
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Global error handlers → always return JSON for /api/* ──
# ──────────────────────────── Admin: User Management ────────────────────────

@app.route("/api/admin/users", methods=["GET"])
@login_required
@admin_required
def admin_list_users():
    if not supabase_admin:
        return jsonify({"success": False, "error": "Supabase not configured"}), 400
    last_error = None
    for attempt in range(5):
        try:
            res = supabase_admin.table("users").select("id,username,tier,display_name,is_active,created_at,last_login").order("created_at", desc=False).execute()
            return jsonify({"success": True, "users": res.data or []})
        except Exception as e:
            last_error = e
            err_str = str(e)
            if ("PGRST002" in err_str or "schema cache" in err_str) and attempt < 4:
                time.sleep(1.0)
            else:
                break
    return jsonify({"success": False, "error": f"Database error: {last_error}"}), 503


@app.route("/api/admin/users", methods=["POST"])
@login_required
@admin_required
def admin_create_user():
    if not supabase_client:
        return jsonify({"success": False, "error": "Supabase not configured"}), 400
    data = request.json
    username = (data.get("username") or "").strip()
    password = (data.get("password") or "").strip()
    tier = data.get("tier", "free")
    display_name = data.get("display_name", username)
    if not username or not password:
        return jsonify({"success": False, "error": "Username and password are required"}), 400
    if len(password) < 4:
        return jsonify({"success": False, "error": "Password must be at least 4 characters"}), 400
    if tier not in ("owner", "admin", "premium", "free"):
        return jsonify({"success": False, "error": "Invalid tier"}), 400
    try:
        # Check if username exists
        existing = supabase_admin.table("users").select("id").eq("username", username).execute()
        if existing.data:
            return jsonify({"success": False, "error": "Username already exists"}), 409
        res = supabase_admin.table("users").insert({
            "username": username,
            "password_hash": _hash_password(password),
            "tier": tier,
            "display_name": display_name,
            "is_active": True,
        }).execute()
        return jsonify({"success": True, "user": res.data[0] if res.data else {}})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/users/<user_id>", methods=["PUT"])
@login_required
@admin_required
def admin_update_user(user_id):
    if not supabase_admin:
        return jsonify({"success": False, "error": "Supabase not configured"}), 400
    data = request.json
    updates = {}
    if "tier" in data and data["tier"] in ("owner", "admin", "premium", "free"):
        updates["tier"] = data["tier"]
    if "display_name" in data:
        updates["display_name"] = data["display_name"]
    if "is_active" in data:
        updates["is_active"] = bool(data["is_active"])
    if "password" in data and data["password"].strip():
        updates["password_hash"] = _hash_password(data["password"].strip())
    if "temp_password" in data:
        tp = data["temp_password"]
        if tp:
            # Setting temp password: also update actual login password and flag for forced change
            updates["temp_password"] = tp
            updates["password_hash"] = _hash_password(tp)
            updates["must_change_password"] = True
        else:
            # Clearing temp password: remove flag too
            updates["temp_password"] = None
            updates["must_change_password"] = False
    if not updates:
        return jsonify({"success": False, "error": "No valid fields to update"}), 400
    try:
        supabase_admin.table("users").update(updates).eq("id", user_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/users/<user_id>", methods=["DELETE"])
@login_required
@admin_required
def admin_delete_user(user_id):
    if not supabase_admin:
        return jsonify({"success": False, "error": "Supabase not configured"}), 400
    # Prevent deleting yourself
    if session.get("user_id") == user_id:
        return jsonify({"success": False, "error": "Cannot delete your own account"}), 400
    try:
        supabase_admin.table("users").delete().eq("id", user_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/user/password", methods=["POST"])
@login_required
def change_own_password():
    if not supabase_client:
        return jsonify({"success": False, "error": "Supabase not configured"}), 400
    data = request.json
    current_pw = (data.get("current_password") or "").strip()
    new_pw = (data.get("new_password") or "").strip()
    if not current_pw or not new_pw:
        return jsonify({"success": False, "error": "Both current and new passwords are required"}), 400
    if len(new_pw) < 4:
        return jsonify({"success": False, "error": "New password must be at least 4 characters"}), 400
    user_id = session.get("user_id")
    if not user_id or user_id == "env_admin":
        return jsonify({"success": False, "error": "Password change only available for Supabase users"}), 400
    try:
        res = supabase_admin.table("users").select("password_hash").eq("id", user_id).execute()
        if not res.data or res.data[0]["password_hash"] != _hash_password(current_pw):
            return jsonify({"success": False, "error": "Current password is incorrect"}), 401
        supabase_admin.table("users").update({
            "password_hash": _hash_password(new_pw),
            "temp_password": None,        # clear temp password
            "must_change_password": False  # clear forced-change flag
        }).eq("id", user_id).execute()
        session["must_change_password"] = False  # update session immediately
        return jsonify({"success": True, "message": "Password changed successfully"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/user/session-info")
@login_required
def session_info():
    """Return user session stats for header display."""
    import requests as http_req
    user_id = session.get("user_id")
    login_time = session.get("login_time", datetime.utcnow().isoformat())
    now = datetime.utcnow()
    # Calculate current session duration
    try:
        lt = datetime.fromisoformat(login_time)
        session_secs = int((now - lt).total_seconds())
    except Exception:
        session_secs = 0
    info = {
        "first_login": None,
        "last_login": None,
        "total_time_minutes": 0,
        "current_session_seconds": session_secs,
        "ip": request.remote_addr or "unknown",
        "location": "—",
    }
    # Get data from Supabase
    if supabase_client and user_id and user_id != "env_admin":
        try:
            res = supabase_admin.table("users").select("created_at,last_login,total_time_minutes").eq("id", user_id).execute()
            if res.data:
                u = res.data[0]
                info["first_login"] = u.get("created_at")
                info["last_login"] = u.get("last_login")
                info["total_time_minutes"] = u.get("total_time_minutes") or 0
        except Exception:
            pass
    elif user_id == "env_admin":
        info["first_login"] = "N/A (env admin)"
        info["last_login"] = login_time
    # Get IP geolocation
    try:
        client_ip = request.headers.get("X-Forwarded-For", request.remote_addr)
        if client_ip:
            client_ip = client_ip.split(",")[0].strip()
            info["ip"] = client_ip
        geo = http_req.get(f"http://ip-api.com/json/{client_ip}?fields=city,regionName,country", timeout=3).json()
        if geo.get("city"):
            info["location"] = f"{geo['city']}, {geo.get('regionName', '')} {geo.get('country', '')}"
    except Exception:
        pass
    info["must_change_password"] = session.get("must_change_password", False)
    return jsonify(info)


# ──────────────────────────── Phase 2: Automation ────────────────────────────
import threading as _threading
import re as _re_auto

SCHEDULER_SECRET = os.getenv("SCHEDULER_SECRET", "")  # set in Render env vars


def _compute_next_run(expr: str, after: datetime = None) -> datetime:
    """Compute next run datetime from a simple schedule expression.
    Formats:
      daily HH:MM          — every day at HH:MM UTC
      weekly DOW HH:MM     — e.g. 'weekly mon 09:00'
      hourly :MM           — every hour at minute MM
      interval N           — every N minutes from now
    """
    if after is None:
        after = datetime.utcnow()
    expr = expr.strip().lower()
    try:
        if expr.startswith("interval "):
            mins = int(expr.split()[1])
            return after + timedelta(minutes=mins)
        if expr.startswith("hourly "):
            minute = int(expr.split(":")[1])
            t = after.replace(second=0, microsecond=0, minute=minute)
            if t <= after:
                t += timedelta(hours=1)
            return t
        if expr.startswith("daily "):
            hm = expr.split()[1]
            h, m = int(hm.split(":")[0]), int(hm.split(":")[1])
            t = after.replace(second=0, microsecond=0, hour=h, minute=m)
            if t <= after:
                t += timedelta(days=1)
            return t
        if expr.startswith("weekly "):
            parts = expr.split()
            day_map = {"mon": 0, "tue": 1, "wed": 2, "thu": 3, "fri": 4, "sat": 5, "sun": 6}
            target_dow = day_map.get(parts[1], 0)
            hm = parts[2]
            h, m = int(hm.split(":")[0]), int(hm.split(":")[1])
            t = after.replace(second=0, microsecond=0, hour=h, minute=m)
            days_ahead = (target_dow - t.weekday()) % 7
            if days_ahead == 0 and t <= after:
                days_ahead = 7
            return t + timedelta(days=days_ahead)
    except Exception:
        pass
    return after + timedelta(hours=1)  # fallback


_sched_lock = _threading.Lock()


def _run_schedule(schedule: dict) -> str:
    """Execute a schedule's prompt via AI and return the result text."""
    try:
        provider = schedule.get("provider", "chatgpt")
        prompt = schedule.get("prompt", "")
        resp = hub.ask(prompt, provider=provider)
        return resp.content if resp.success else f"[Error] {resp.error}"
    except Exception as e:
        return f"[Error] {e}"


def _process_due_schedules(username: str = None):
    """Find and run any schedules whose next_run_at has passed."""
    if not supabase_client:
        return []
    fired = []
    try:
        now_iso = datetime.utcnow().isoformat()
        q = supabase_client.table("schedules").select("*").eq("is_active", True).lte("next_run_at", now_iso)
        if username:
            q = q.eq("user_id", username)
        result = q.execute()
        for sched in (result.data or []):
            if _sched_lock.locked():
                continue
            with _sched_lock:
                # Re-check to avoid double-fire across workers
                fresh = supabase_client.table("schedules").select("next_run_at,is_running").eq("id", sched["id"]).execute()
                if not fresh.data or fresh.data[0].get("is_running"):
                    continue
                supabase_client.table("schedules").update({"is_running": True}).eq("id", sched["id"]).execute()
            try:
                result_text = _run_schedule(sched)
                next_run = _compute_next_run(sched["schedule_expr"])
                update = {
                    "last_run_at": datetime.utcnow().isoformat(),
                    "next_run_at": next_run.isoformat(),
                    "last_result": result_text[:2000],
                    "is_running": False,
                }
                # Save to workspace if folder_id set
                if sched.get("folder_id") and supabase_client:
                    try:
                        fname = f"{sched['name']}_{datetime.utcnow().strftime('%Y%m%d_%H%M')}.txt"
                        supabase_client.table("workspace_files").insert({
                            "folder_id": sched["folder_id"],
                            "user_id": sched["user_id"],
                            "name": fname,
                            "type": "note",
                            "content": result_text,
                            "metadata": {"source": "scheduler", "schedule_id": sched["id"]},
                        }).execute()
                    except Exception:
                        pass
                supabase_client.table("schedules").update(update).eq("id", sched["id"]).execute()
                fired.append({"id": sched["id"], "name": sched["name"]})
            except Exception as e:
                supabase_client.table("schedules").update({"is_running": False}).eq("id", sched["id"]).execute()
    except Exception:
        pass
    return fired


# Start APScheduler (fires _process_due_schedules every 60s)
try:
    from apscheduler.schedulers.background import BackgroundScheduler
    _bg_scheduler = BackgroundScheduler(daemon=True)
    _bg_scheduler.add_job(_process_due_schedules, "interval", seconds=60, max_instances=1, coalesce=True)
    _bg_scheduler.start()
    print("  ✅ APScheduler started (60s interval)")
except Exception as _e:
    print(f"  ⚠️ APScheduler not available: {_e}")
    _bg_scheduler = None


# ── Schedule APIs ──────────────────────────────────────────────────────────

@app.route("/api/schedules", methods=["GET"])
@login_required
def api_schedules_list():
    if not supabase_client:
        return jsonify({"schedules": []})
    try:
        r = supabase_client.table("schedules").select("*").eq(
            "user_id", session.get("username", "admin")
        ).order("created_at", desc=False).execute()
        return jsonify({"schedules": r.data or []})
    except Exception as e:
        return jsonify({"schedules": [], "error": str(e)})


@app.route("/api/schedules", methods=["POST"])
@login_required
def api_schedules_create():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    expr = data.get("schedule_expr", "daily 09:00")
    next_run = _compute_next_run(expr)
    try:
        r = supabase_client.table("schedules").insert({
            "user_id": session.get("username", "admin"),
            "name": data.get("name", "New Schedule"),
            "schedule_expr": expr,
            "prompt": data.get("prompt", ""),
            "provider": data.get("provider", "chatgpt"),
            "folder_id": data.get("folder_id") or None,
            "is_active": True,
            "is_running": False,
            "next_run_at": next_run.isoformat(),
        }).execute()
        return jsonify({"schedule": r.data[0] if r.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/schedules/<sched_id>", methods=["PUT"])
@login_required
def api_schedules_update(sched_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    update = {}
    for f in ["name", "prompt", "provider", "folder_id", "is_active"]:
        if f in data:
            update[f] = data[f]
    if "schedule_expr" in data:
        update["schedule_expr"] = data["schedule_expr"]
        update["next_run_at"] = _compute_next_run(data["schedule_expr"]).isoformat()
    if not update:
        return jsonify({"success": True})
    try:
        supabase_client.table("schedules").update(update).eq("id", sched_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/schedules/<sched_id>", methods=["DELETE"])
@login_required
def api_schedules_delete(sched_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("schedules").delete().eq("id", sched_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/schedules/<sched_id>/run", methods=["POST"])
@login_required
def api_schedules_run(sched_id):
    """Manually trigger a schedule immediately."""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        r = supabase_client.table("schedules").select("*").eq("id", sched_id).execute()
        if not r.data:
            return jsonify({"error": "Not found"}), 404
        sched = r.data[0]
        result_text = _run_schedule(sched)
        next_run = _compute_next_run(sched["schedule_expr"])
        supabase_client.table("schedules").update({
            "last_run_at": datetime.utcnow().isoformat(),
            "next_run_at": next_run.isoformat(),
            "last_result": result_text[:2000],
            "is_running": False,
        }).eq("id", sched_id).execute()
        if sched.get("folder_id"):
            try:
                fname = f"{sched['name']}_{datetime.utcnow().strftime('%Y%m%d_%H%M')}.txt"
                supabase_client.table("workspace_files").insert({
                    "folder_id": sched["folder_id"],
                    "user_id": sched["user_id"],
                    "name": fname, "type": "note",
                    "content": result_text,
                    "metadata": {"source": "scheduler", "schedule_id": sched_id},
                }).execute()
            except Exception:
                pass
        return jsonify({"success": True, "result": result_text[:500]})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/schedules/check")
def api_schedules_check():
    """External pinger endpoint (UptimeRobot/GitHub Actions).
    Protected by SCHEDULER_SECRET env var if set.
    """
    if SCHEDULER_SECRET:
        token = request.args.get("secret") or request.headers.get("X-Scheduler-Secret", "")
        if token != SCHEDULER_SECRET:
            return jsonify({"error": "Unauthorized"}), 401
    fired = _process_due_schedules()
    return jsonify({"fired": fired, "count": len(fired), "ts": datetime.utcnow().isoformat()})


# ── Workflow APIs ────────────────────────────────────────────────────────────

@app.route("/api/workflows", methods=["GET"])
@login_required
def api_workflows_list():
    if not supabase_client:
        return jsonify({"workflows": []})
    try:
        r = supabase_client.table("workflows").select("id,name,description,steps,created_at").eq(
            "user_id", session.get("username", "admin")
        ).order("created_at", desc=False).execute()
        return jsonify({"workflows": r.data or []})
    except Exception as e:
        return jsonify({"workflows": [], "error": str(e)})


@app.route("/api/workflows", methods=["POST"])
@login_required
def api_workflows_create():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    try:
        r = supabase_client.table("workflows").insert({
            "user_id": session.get("username", "admin"),
            "name": data.get("name", "New Workflow"),
            "description": data.get("description", ""),
            "steps": data.get("steps", []),
        }).execute()
        return jsonify({"workflow": r.data[0] if r.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/workflows/<wf_id>", methods=["PUT"])
@login_required
def api_workflows_update(wf_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    update = {k: data[k] for k in ["name", "description", "steps"] if k in data}
    try:
        supabase_client.table("workflows").update(update).eq("id", wf_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/workflows/<wf_id>", methods=["DELETE"])
@login_required
def api_workflows_delete(wf_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("workflows").delete().eq("id", wf_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/workflows/<wf_id>/run", methods=["POST"])
@login_required
def api_workflows_run(wf_id):
    """Execute workflow steps sequentially. Returns SSE stream of progress."""
    if not supabase_client:
        return jsonify({"error": "No database"}), 400

    data = request.json or {}
    user_input = data.get("input", "")
    folder_id = data.get("folder_id")

    try:
        r = supabase_client.table("workflows").select("*").eq("id", wf_id).execute()
        if not r.data:
            return jsonify({"error": "Not found"}), 404
        wf = r.data[0]
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    steps = wf.get("steps") or []
    if not steps:
        return jsonify({"error": "워크플로우에 단계가 없습니다."}), 400

    def generate():
        prev_output = user_input
        all_outputs = []
        for i, step in enumerate(steps):
            step_name = step.get("name", f"Step {i+1}")
            provider = step.get("provider", "chatgpt")
            prompt_tmpl = step.get("prompt", "")
            # Template substitution
            prompt = prompt_tmpl.replace("{{input}}", user_input).replace("{{prev_output}}", prev_output)
            yield f"data: {json.dumps({'type':'step_start','step':i+1,'total':len(steps),'name':step_name})}\n\n"
            try:
                resp = hub.ask(prompt, provider=provider)
                output = resp.content if resp.success else f"[Error] {resp.error}"
            except Exception as ex:
                output = f"[Error] {ex}"
            prev_output = output
            all_outputs.append({"step": i+1, "name": step_name, "provider": provider, "output": output})
            yield f"data: {json.dumps({'type':'step_done','step':i+1,'name':step_name,'output':output[:1000]})}\n\n"
        # Save final result to workspace
        if folder_id:
            try:
                full_report = "\n\n".join([
                    f"## Step {o['step']}: {o['name']} ({o['provider']})\n{o['output']}"
                    for o in all_outputs
                ])
                fname = f"{wf['name']}_{datetime.utcnow().strftime('%Y%m%d_%H%M')}.txt"
                supabase_client.table("workspace_files").insert({
                    "folder_id": folder_id,
                    "user_id": session.get("username", "admin"),
                    "name": fname, "type": "note",
                    "content": full_report,
                    "metadata": {"source": "workflow", "workflow_id": wf_id},
                }).execute()
            except Exception:
                pass
        yield f"data: {json.dumps({'type':'done','steps':len(steps),'final':prev_output[:2000]})}\n\n"

    return Response(stream_with_context(generate()), mimetype="text/event-stream",
                    headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


# ── Webhook APIs ──────────────────────────────────────────────────────────────

@app.route("/api/webhooks", methods=["GET"])
@login_required
def api_webhooks_list():
    if not supabase_client:
        return jsonify({"webhooks": []})
    try:
        r = supabase_client.table("webhooks").select("id,name,token,prompt_template,provider,is_active,last_triggered_at,created_at").eq(
            "user_id", session.get("username", "admin")
        ).order("created_at", desc=False).execute()
        return jsonify({"webhooks": r.data or []})
    except Exception as e:
        return jsonify({"webhooks": [], "error": str(e)})


@app.route("/api/webhooks", methods=["POST"])
@login_required
def api_webhooks_create():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    token = secrets.token_hex(20)
    try:
        r = supabase_client.table("webhooks").insert({
            "user_id": session.get("username", "admin"),
            "name": data.get("name", "New Webhook"),
            "token": token,
            "prompt_template": data.get("prompt_template", "Summarize this payload: {{payload}}"),
            "provider": data.get("provider", "chatgpt"),
            "folder_id": data.get("folder_id") or None,
            "is_active": True,
        }).execute()
        return jsonify({"webhook": r.data[0] if r.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/webhooks/<wh_id>", methods=["PUT"])
@login_required
def api_webhooks_update(wh_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    update = {k: data[k] for k in ["name", "prompt_template", "provider", "is_active", "folder_id"] if k in data}
    try:
        supabase_client.table("webhooks").update(update).eq("id", wh_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/webhooks/<wh_id>", methods=["DELETE"])
@login_required
def api_webhooks_delete(wh_id):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("webhooks").delete().eq("id", wh_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/webhook/<token>", methods=["POST"])
def api_webhook_receive(token):
    """Public webhook endpoint — no login required, token-protected."""
    if not supabase_client:
        return jsonify({"error": "No database"}), 503
    try:
        r = supabase_client.table("webhooks").select("*").eq("token", token).eq("is_active", True).execute()
        if not r.data:
            return jsonify({"error": "Webhook not found or inactive"}), 404
        wh = r.data[0]
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    payload = request.json or request.form.to_dict() or {}
    payload_str = json.dumps(payload, ensure_ascii=False)

    # Template substitution: {{payload}} → full JSON, {{payload.key}} → field
    def _sub(m):
        key = m.group(1).strip()
        if key == "payload":
            return payload_str
        if key.startswith("payload."):
            field = key[8:]
            return str(payload.get(field, ""))
        return m.group(0)

    prompt = _re_auto.sub(r"\{\{([^}]+)\}\}", _sub, wh.get("prompt_template", ""))

    try:
        resp = hub.ask(prompt, provider=wh.get("provider", "chatgpt"))
        result_text = resp.content if resp.success else f"[Error] {resp.error}"
        # Save to workspace
        if wh.get("folder_id"):
            try:
                supabase_client.table("workspace_files").insert({
                    "folder_id": wh["folder_id"],
                    "user_id": wh["user_id"],
                    "name": f"webhook_{wh['name']}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.txt",
                    "type": "note",
                    "content": result_text,
                    "metadata": {"source": "webhook", "webhook_id": wh["id"], "payload": payload},
                }).execute()
            except Exception:
                pass
        supabase_client.table("webhooks").update({
            "last_triggered_at": datetime.utcnow().isoformat()
        }).eq("id", wh["id"]).execute()
        return jsonify({"success": True, "result": result_text[:1000]})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ──────────────────────────── Phase 3: Integrations ────────────────────────────
import smtplib as _smtplib
from email.mime.text import MIMEText as _MIMEText
from email.mime.multipart import MIMEMultipart as _MIMEMultipart
import http.client as _http_client
import urllib.parse as _urlparse


def _integrations_for(user: str, itype: str = None) -> list:
    """Fetch active integrations from DB."""
    if not supabase_client:
        return []
    try:
        q = supabase_client.table("integrations").select("*").eq("user_id", user).eq("is_active", True)
        if itype:
            q = q.eq("type", itype)
        return q.execute().data or []
    except Exception:
        return []


def _send_slack(webhook_url: str, text: str) -> tuple[bool, str]:
    """Post message to Slack incoming webhook."""
    try:
        import requests as _req
        r = _req.post(webhook_url, json={"text": text}, timeout=10)
        if r.status_code == 200:
            return True, "ok"
        return False, f"HTTP {r.status_code}: {r.text[:200]}"
    except Exception as e:
        return False, str(e)


def _send_notion_page(token: str, database_id: str, title: str, content: str) -> tuple[bool, str]:
    """Create a page in a Notion database."""
    try:
        import requests as _req
        body = {
            "parent": {"database_id": database_id},
            "properties": {
                "Name": {"title": [{"text": {"content": title[:100]}}]}
            },
            "children": [
                {"object": "block", "type": "paragraph",
                 "paragraph": {"rich_text": [{"text": {"content": seg}}]}}
                for seg in [content[i:i+2000] for i in range(0, min(len(content), 10000), 2000)]
            ]
        }
        r = _req.post(
            "https://api.notion.com/v1/pages",
            headers={"Authorization": f"Bearer {token}", "Notion-Version": "2022-06-28",
                     "Content-Type": "application/json"},
            json=body, timeout=15
        )
        if r.status_code in (200, 201):
            return True, r.json().get("url", "")
        return False, f"HTTP {r.status_code}: {r.text[:300]}"
    except Exception as e:
        return False, str(e)


def _send_email(cfg: dict, subject: str, body: str) -> tuple[bool, str]:
    """Send email via SMTP (Gmail app password or any SMTP)."""
    try:
        msg = _MIMEMultipart()
        msg["From"] = cfg["user"]
        msg["To"] = cfg.get("to", cfg["user"])
        msg["Subject"] = subject
        msg.attach(_MIMEText(body, "plain", "utf-8"))
        with _smtplib.SMTP(cfg.get("host", "smtp.gmail.com"), int(cfg.get("port", 587))) as smtp:
            smtp.ehlo()
            smtp.starttls()
            smtp.login(cfg["user"], cfg["password"])
            smtp.sendmail(cfg["user"], msg["To"], msg.as_string())
        return True, f"Sent to {msg['To']}"
    except Exception as e:
        return False, str(e)


def _get_calendar_events(ical_url: str, max_events: int = 10) -> str:
    """Fetch upcoming events from a Google Calendar iCal URL."""
    try:
        import requests as _req
        from datetime import timezone
        r = _req.get(ical_url, timeout=10)
        if r.status_code != 200:
            return f"[Calendar fetch error: HTTP {r.status_code}]"
        text = r.text
        now = datetime.now(timezone.utc)
        events = []
        # Simple iCal parser — no external library needed
        for vevent in text.split("BEGIN:VEVENT")[1:]:
            try:
                def _get(field):
                    for line in vevent.splitlines():
                        if line.startswith(field + ":") or line.startswith(field + ";"):
                            return line.split(":", 1)[1].strip() if ":" in line else ""
                    return ""
                summary = _get("SUMMARY")
                dtstart_raw = _get("DTSTART")
                dtend_raw = _get("DTEND")
                # Parse date
                def _parse_dt(s):
                    s = s.replace("Z", "").replace("-", "").replace(":", "").replace("T", "")
                    if len(s) >= 8:
                        try:
                            return datetime(int(s[:4]), int(s[4:6]), int(s[6:8]),
                                            int(s[8:10]) if len(s) > 8 else 0,
                                            int(s[10:12]) if len(s) > 10 else 0,
                                            tzinfo=timezone.utc)
                        except Exception:
                            return None
                    return None
                start = _parse_dt(dtstart_raw)
                if start and start >= now:
                    events.append((start, summary, dtend_raw))
            except Exception:
                continue
        events.sort(key=lambda x: x[0])
        if not events:
            return "일정 없음"
        lines = [f"• {e[0].strftime('%m/%d %H:%M')} {e[1]}" for e in events[:max_events]]
        return "\n".join(lines)
    except Exception as e:
        return f"[Calendar error: {e}]"


# ── Group APIs ───────────────────────────────────────────────────────────────

def _user_group_ids(username: str) -> list:
    """Return list of group IDs the user belongs to."""
    if not supabase_client:
        return []
    try:
        r = supabase_client.table("group_members").select("group_id").eq("user_id", username).execute()
        return [row["group_id"] for row in (r.data or [])]
    except Exception:
        return []


def _is_admin(username: str) -> bool:
    """Check if user is system admin (env admin or tier=admin/owner in DB)."""
    if username == APP_USERNAME:
        return True
    try:
        r = supabase_client.table("users").select("tier").eq("username", username).execute()
        return bool(r.data and r.data[0].get("tier") in ("admin", "owner"))
    except Exception:
        return False


def _is_owner(username: str) -> bool:
    """Check if user is owner tier only."""
    if username == APP_USERNAME:
        return True
    try:
        r = supabase_client.table("users").select("tier").eq("username", username).execute()
        return bool(r.data and r.data[0].get("tier") == "owner")
    except Exception:
        return False


@app.route("/api/groups", methods=["GET"])
@login_required
def api_groups_list():
    """List all groups (owner only)."""
    if not supabase_client:
        return jsonify({"groups": []})
    username = session.get("username", "admin")
    if not _is_owner(username):
        return jsonify({"groups": [], "error": "owner 전용 기능입니다."}), 403
    try:
        r = supabase_client.table("groups").select("*").order("created_at").execute()
        groups = r.data or []
        # Attach member counts
        for g in groups:
            try:
                mc = supabase_client.table("group_members").select("id").eq("group_id", g["id"]).execute()
                g["member_count"] = len(mc.data or [])
            except Exception:
                g["member_count"] = 0
        return jsonify({"groups": groups})
    except Exception as e:
        return jsonify({"groups": [], "error": str(e)})


@app.route("/api/groups", methods=["POST"])
@login_required
def api_groups_create():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    username = session.get("username", "admin")
    if not _is_owner(username):
        return jsonify({"error": "owner만 그룹을 생성할 수 있습니다."}), 403
    data = request.json or {}
    try:
        r = supabase_client.table("groups").insert({
            "name": data.get("name", "New Group"),
            "description": data.get("description", ""),
            "created_by": username,
        }).execute()
        group = r.data[0] if r.data else {}
        # Auto-add creator as admin member
        if group.get("id"):
            supabase_client.table("group_members").insert({
                "group_id": group["id"], "user_id": username, "role": "admin"
            }).execute()
        return jsonify({"group": group, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<gid>", methods=["PUT"])
@login_required
def api_groups_update(gid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    username = session.get("username", "admin")
    if not _is_owner(username):
        return jsonify({"error": "권한 없음"}), 403
    data = request.json or {}
    update = {k: data[k] for k in ["name", "description"] if k in data}
    try:
        supabase_client.table("groups").update(update).eq("id", gid).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<gid>", methods=["DELETE"])
@login_required
def api_groups_delete(gid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    if not _is_owner(session.get("username", "admin")):
        return jsonify({"error": "권한 없음"}), 403
    try:
        supabase_client.table("group_members").delete().eq("group_id", gid).execute()
        supabase_client.table("groups").delete().eq("id", gid).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<gid>/members", methods=["GET"])
@login_required
def api_group_members_list(gid):
    if not supabase_client:
        return jsonify({"members": []})
    try:
        r = supabase_client.table("group_members").select("*").eq("group_id", gid).order("created_at").execute()
        return jsonify({"members": r.data or []})
    except Exception as e:
        return jsonify({"members": [], "error": str(e)})


@app.route("/api/groups/<gid>/members", methods=["POST"])
@login_required
def api_group_members_add(gid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    if not _is_owner(session.get("username", "admin")):
        return jsonify({"error": "권한 없음"}), 403
    data = request.json or {}
    user_id = data.get("user_id", "").strip()
    if not user_id:
        return jsonify({"error": "user_id 필요"}), 400
    try:
        # Check if member already exists, then update or insert
        existing = supabase_client.table("group_members").select("id").eq(
            "group_id", gid).eq("user_id", user_id).execute()
        if existing.data:
            supabase_client.table("group_members").update(
                {"role": data.get("role", "member")}
            ).eq("id", existing.data[0]["id"]).execute()
        else:
            supabase_client.table("group_members").insert({
                "group_id": gid,
                "user_id": user_id,
                "role": data.get("role", "member"),
            }).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<gid>/members/<uid>", methods=["DELETE"])
@login_required
def api_group_members_remove(gid, uid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    if not _is_owner(session.get("username", "admin")):
        return jsonify({"error": "권한 없음"}), 403
    try:
        supabase_client.table("group_members").delete().eq("group_id", gid).eq("user_id", uid).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ── Integration CRUD ──────────────────────────────────────────────────────────

@app.route("/api/integrations", methods=["GET"])
@login_required
def api_integrations_list():
    """List integrations visible to current user:
    personal/null (own) + group (member of) + global.
    """
    if not supabase_client:
        return jsonify({"integrations": []})
    username = session.get("username", "admin")
    try:
        results = []
        seen_ids = set()
        # 1. Personal — includes rows where scope IS NULL (legacy rows)
        r = supabase_client.table("integrations").select(
            "id,type,name,scope,group_id,is_active,last_used_at,created_at,user_id"
        ).eq("user_id", username).execute()
        for row in (r.data or []):
            scope = row.get("scope") or "personal"
            if scope == "personal":
                row["_label"] = "개인"
                results.append(row); seen_ids.add(row["id"])
        # 2. Group integrations
        gids = _user_group_ids(username)
        if gids:
            try:
                gr = supabase_client.table("integrations").select(
                    "id,type,name,scope,group_id,is_active,last_used_at,created_at,user_id"
                ).eq("scope", "group").in_("group_id", gids).execute()
                for row in (gr.data or []):
                    if row["id"] not in seen_ids:
                        try:
                            gn = supabase_client.table("groups").select("name").eq("id", row["group_id"]).execute()
                            row["_label"] = gn.data[0]["name"] if gn.data else "그룹"
                        except Exception:
                            row["_label"] = "그룹"
                        results.append(row); seen_ids.add(row["id"])
            except Exception:
                pass
        # 3. Global
        try:
            glob_r = supabase_client.table("integrations").select(
                "id,type,name,scope,group_id,is_active,last_used_at,created_at,user_id"
            ).eq("scope", "global").execute()
            for row in (glob_r.data or []):
                if row["id"] not in seen_ids:
                    row["_label"] = "전체공유"
                    results.append(row); seen_ids.add(row["id"])
        except Exception:
            pass
        return jsonify({"integrations": results})
    except Exception as e:
        return jsonify({"integrations": [], "error": str(e)})


@app.route("/api/integrations", methods=["POST"])
@login_required
def api_integrations_create():
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    data = request.json or {}
    try:
        # Upsert by type+name
        existing = supabase_client.table("integrations").select("id").eq(
            "user_id", session.get("username", "admin")
        ).eq("type", data.get("type", "")).eq("name", data.get("name", "")).execute()
        scope = data.get("scope", "personal")
        group_id = data.get("group_id") or None
        # Only admin can set global scope
        if scope == "global" and not _is_admin(session.get("username", "admin")):
            scope = "group"
        payload = {
            "user_id": session.get("username", "admin"),
            "type": data.get("type"),
            "name": data.get("name", data.get("type", "Integration")),
            "config": data.get("config", {}),
            "scope": scope,
            "group_id": group_id,
            "is_active": True,
        }
        if existing.data:
            supabase_client.table("integrations").update(payload).eq("id", existing.data[0]["id"]).execute()
            iid = existing.data[0]["id"]
        else:
            r = supabase_client.table("integrations").insert(payload).execute()
            iid = r.data[0]["id"] if r.data else None
        return jsonify({"id": iid, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/integrations/<iid>", methods=["DELETE"])
@login_required
def api_integrations_delete(iid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        supabase_client.table("integrations").delete().eq("id", iid).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/integrations/<iid>/test", methods=["POST"])
@login_required
def api_integrations_test(iid):
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        r = supabase_client.table("integrations").select("*").eq("id", iid).execute()
        if not r.data:
            return jsonify({"error": "Not found"}), 404
        integ = r.data[0]
        cfg = integ.get("config") or {}
        itype = integ.get("type")
        test_msg = "✅ AI Hub 연동 테스트 메시지입니다."
        if itype == "slack":
            ok, msg = _send_slack(cfg.get("webhook_url", ""), test_msg)
        elif itype == "notion":
            ok, msg = _send_notion_page(cfg.get("token", ""), cfg.get("database_id", ""), "AI Hub 테스트", test_msg)
        elif itype == "email":
            ok, msg = _send_email(cfg, "AI Hub 테스트 이메일", test_msg)
        elif itype == "calendar":
            events = _get_calendar_events(cfg.get("ical_url", ""))
            ok, msg = True, f"일정 미리보기:\n{events}"
        else:
            ok, msg = False, "Unknown integration type"
        return jsonify({"success": ok, "message": msg})
    except Exception as e:
        return jsonify({"success": False, "message": str(e)})


@app.route("/api/integrations/send", methods=["POST"])
@login_required
def api_integrations_send():
    """Send text to one or more integrations by type."""
    data = request.json or {}
    itype = data.get("type")
    text = data.get("text", "")
    title = data.get("title", "AI Hub 결과")
    username = session.get("username", "admin")
    integ_id = data.get("integration_id")
    if not supabase_client:
        return jsonify({"error": "No database"}), 400
    try:
        if integ_id:
            r = supabase_client.table("integrations").select("*").eq("id", integ_id).execute()
        else:
            r = supabase_client.table("integrations").select("*").eq("user_id", username).eq("type", itype).eq("is_active", True).execute()
        integs = r.data or []
        if not integs:
            return jsonify({"error": "연동 설정이 없습니다."}), 400
        results = []
        for integ in integs:
            cfg = integ.get("config") or {}
            t = integ.get("type")
            if t == "slack":
                ok, msg = _send_slack(cfg.get("webhook_url", ""), f"*{title}*\n{text}")
            elif t == "notion":
                ok, msg = _send_notion_page(cfg.get("token", ""), cfg.get("database_id", ""), title, text)
            elif t == "email":
                ok, msg = _send_email(cfg, title, text)
            else:
                ok, msg = False, f"Send not supported for type: {t}"
            if ok:
                supabase_client.table("integrations").update({"last_used_at": datetime.utcnow().isoformat()}).eq("id", integ["id"]).execute()
            results.append({"id": integ["id"], "name": integ["name"], "success": ok, "message": msg})
        all_ok = all(r["success"] for r in results)
        return jsonify({"success": all_ok, "results": results})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/integrations/calendar/events", methods=["GET"])
@login_required
def api_calendar_events():
    """Return upcoming calendar events as text (for prompt injection)."""
    username = session.get("username", "admin")
    integ_id = request.args.get("id")
    if not supabase_client:
        return jsonify({"events": "", "error": "No database"})
    try:
        if integ_id:
            r = supabase_client.table("integrations").select("*").eq("id", integ_id).execute()
        else:
            r = supabase_client.table("integrations").select("*").eq("user_id", username).eq("type", "calendar").eq("is_active", True).execute()
        integs = r.data or []
        if not integs:
            return jsonify({"events": "캘린더 연동이 설정되지 않았습니다.", "count": 0})
        cfg = integs[0].get("config") or {}
        events = _get_calendar_events(cfg.get("ical_url", ""), max_events=15)
        return jsonify({"events": events, "count": len(events.splitlines())})
    except Exception as e:
        return jsonify({"events": "", "error": str(e)})


# ── Favicon & common browser requests ──
@app.route("/favicon.ico")
def favicon():
    return "", 204  # No Content — silently ignore

@app.route("/robots.txt")
def robots():
    return "User-agent: *\nDisallow: /api/\n", 200, {"Content-Type": "text/plain"}

# ──────────────────────────── Error Handlers ────────────────────────────

@app.errorhandler(Exception)
def handle_exception(e):
    import traceback
    if request.path.startswith("/api/"):
        return jsonify({"success": False, "error": str(e),
                        "detail": traceback.format_exc()[-500:]}), 500
    raise e


@app.errorhandler(500)
def handle_500(e):
    if request.path.startswith("/api/"):
        return jsonify({"success": False, "error": "Internal server error", "detail": str(e)}), 500
    return str(e), 500


@app.route("/api/account/change-password", methods=["POST"])
@login_required
def api_change_password():
    username = session.get("username", "")
    data = request.json or {}
    current_pw = data.get("current_password", "")
    new_pw = data.get("new_password", "")
    if not current_pw or not new_pw:
        return jsonify({"error": "필드 누락"}), 400
    if len(new_pw) < 6:
        return jsonify({"error": "새 비밀번호는 6자 이상이어야 합니다."}), 400
    # Verify current password
    current_hash = _hash_password(current_pw)
    if not supabase_client:
        # Fallback: env-based user
        if username == APP_USERNAME and current_hash == APP_PASSWORD_HASH:
            return jsonify({"error": "환경변수 기반 계정은 Render 대시보드에서 변경하세요."}), 400
        return jsonify({"error": "DB 연결 없음"}), 400
    try:
        r = supabase_client.table("users").select("id,password_hash").eq("username", username).execute()
        if not r.data:
            return jsonify({"error": "유저를 찾을 수 없습니다."}), 404
        user = r.data[0]
        if user["password_hash"] != current_hash:
            return jsonify({"error": "현재 비밀번호가 올바르지 않습니다."}), 403
        new_hash = _hash_password(new_pw)
        supabase_client.table("users").update({"password_hash": new_hash}).eq("id", user["id"]).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/admin/temp-password", methods=["POST"])
@login_required
def api_temp_password():
    if not _is_owner(session.get("username", "")):
        return jsonify({"error": "owner 전용 기능입니다."}), 403
    if not supabase_client:
        return jsonify({"error": "DB 연결 없음"}), 400
    data = request.json or {}
    user_id = data.get("user_id", "")
    custom_pw = data.get("password", "").strip()
    if not user_id:
        return jsonify({"error": "user_id 필요"}), 400
    if not custom_pw or len(custom_pw) < 6:
        return jsonify({"error": "비밀번호는 6자 이상이어야 합니다."}), 400
    try:
        new_hash = _hash_password(custom_pw)
        supabase_client.table("users").update({"password_hash": new_hash}).eq("id", user_id).execute()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/user-templates", methods=["GET"])
@login_required
def api_user_templates_list():
    if not supabase_client:
        return jsonify({"templates": []})
    user_id = session.get("username", "")
    ttype = request.args.get("type", "")
    try:
        q = supabase_client.table("user_templates").select("*").eq("user_id", user_id).order("created_at")
        if ttype:
            q = q.eq("type", ttype)
        r = q.execute()
        return jsonify({"templates": r.data or []})
    except Exception as e:
        return jsonify({"templates": [], "error": str(e)})


@app.route("/api/user-templates", methods=["POST"])
@login_required
def api_user_templates_create():
    if not supabase_client:
        return jsonify({"error": "DB 연결 없음"}), 400
    user_id = session.get("username", "")
    data = request.json or {}
    name = data.get("name", "").strip()
    prompt = data.get("prompt", "").strip()
    ttype = data.get("type", "template")
    if not name or not prompt:
        return jsonify({"error": "이름과 프롬프트를 입력하세요."}), 400
    if ttype not in ("template", "strategy"):
        return jsonify({"error": "type은 template 또는 strategy"}), 400
    try:
        r = supabase_client.table("user_templates").insert({
            "user_id": user_id, "type": ttype, "name": name, "prompt": prompt
        }).execute()
        return jsonify({"template": r.data[0] if r.data else {}, "success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/user-templates/<tid>", methods=["DELETE"])
@login_required
def api_user_templates_delete(tid):
    if not supabase_client:
        return jsonify({"error": "DB 연결 없음"}), 400
    user_id = session.get("username", "")
    try:
        supabase_client.table("user_templates").delete().eq("id", tid).eq("user_id", user_id).execute()
        return jsonify({"deleted": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/youtube/transcript", methods=["POST"])
@login_required
def api_youtube_transcript():
    data = request.json or {}
    url = data.get("url", "").strip()
    if not url:
        return jsonify({"error": "URL이 필요합니다."}), 400
    import re
    m = re.search(r"(?:youtube\.com/watch\?v=|youtu\.be/|youtube\.com/embed/|youtube\.com/live/|youtube\.com/shorts/)([a-zA-Z0-9_-]{11})", url)
    if not m:
        return jsonify({"error": "올바른 YouTube URL이 아닙니다."}), 400
    vid = m.group(1)

    def _get_text(snippet):
        """dict 또는 객체 모두 처리 (youtube-transcript-api 버전 호환)"""
        if isinstance(snippet, dict):
            return snippet.get("text", "")
        return getattr(snippet, "text", str(snippet))

    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        transcript = None
        lang_used = ""

        try:
            tlist = YouTubeTranscriptApi.list_transcripts(vid)
            # 1) 수동 자막: ko → en → 기타
            for lang in ["ko", "en", "en-US", "en-GB"]:
                try:
                    transcript = tlist.find_transcript([lang]).fetch()
                    lang_used = lang
                    break
                except Exception:
                    pass
            # 2) 자동 생성 자막
            if transcript is None:
                try:
                    auto = tlist.find_generated_transcript(["ko", "en", "en-US"])
                    transcript = auto.fetch()
                    lang_used = "auto"
                except Exception:
                    pass
            # 3) 첫 번째 사용 가능한 자막
            if transcript is None:
                for t in tlist:
                    transcript = t.fetch()
                    lang_used = t.language_code
                    break
        except Exception:
            pass

        # 4) 최후 수단
        if transcript is None:
            transcript = YouTubeTranscriptApi.get_transcript(vid)
            lang_used = "default"

        text = " ".join([_get_text(s) for s in transcript if _get_text(s).strip()])

        if not text.strip():
            return jsonify({"error": "자막 내용이 비어 있습니다."}), 400

        if len(text) > 14000:
            text = text[:14000] + "\n...(이하 생략)"

        return jsonify({"transcript": text, "video_id": vid, "lang": lang_used, "success": True})

    except Exception as e:
        err = str(e)
        if "TranscriptsDisabled" in err or "disabled" in err.lower():
            msg = "이 영상은 자막이 비활성화되어 있습니다."
        elif "NoTranscriptFound" in err or "Could not find" in err:
            msg = "사용 가능한 자막이 없습니다. (자막 없는 영상)"
        elif "VideoUnavailable" in err:
            msg = "영상을 찾을 수 없습니다. (비공개 또는 삭제됨)"
        else:
            msg = f"자막 오류: {err}"
        return jsonify({"error": msg}), 400


if __name__ == "__main__":
    _seed_admin_user()
    port = int(os.getenv("PORT", 5000))
    print(f"\n  AI Hub Web App")
    print(f"  http://localhost:{port}")
    print(f"  Login: {APP_USERNAME}\n")
    app.run(debug=False, host="0.0.0.0", port=port)
else:
    # Also run seed when imported by gunicorn
    _seed_admin_user()
